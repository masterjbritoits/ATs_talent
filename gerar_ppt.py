#!/usr/bin/env python3
"""
Gera apresentação comercial premium para ITSector Talent Inbox ATS
Design: Dark premium com acentos azul/ciano, tipografia moderna
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ─── CORES PREMIUM ───────────────────────────────────────────────
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)        # Fundo escuro principal
DARK_BG2 = RGBColor(0x15, 0x20, 0x3C)       # Fundo escuro secundário
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)    # Azul principal
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)    # Ciano destaque
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)   # Verde sucesso
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)  # Roxo destaque
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)   # Amber/Dourado
ACCENT_ROSE = RGBColor(0xF4, 0x3F, 0x5E)    # Rosa/Vermelho
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)     # Texto secundário
MED_GRAY = RGBColor(0x94, 0xA3, 0xB8)       # Texto terciário
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)        # Fundo de cartão
GRADIENT_START = RGBColor(0x1E, 0x40, 0xAF)  # Gradiente azul
GRADIENT_END = RGBColor(0x06, 0xB6, 0xD4)    # Gradiente ciano

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# ─── HELPERS ─────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rgb_tuple(c):
    return (c[0], c[1], c[2])

def darken(c, factor=4):
    r, g, b = rgb_tuple(c)
    return RGBColor(r // factor, g // factor, b // factor)

def add_gradient_bar(slide, left, top, width, height, color1, color2):
    """Adiciona barra com gradiente simulado (3 segmentos)"""
    seg_w = width // 3
    r1, g1, b1 = rgb_tuple(color1)
    r2, g2, b2 = rgb_tuple(color2)
    colors = [color1, RGBColor(
        (r1 + r2) // 2,
        (g1 + g2) // 2,
        (b1 + b2) // 2
    ), color2]
    for i, c in enumerate(colors):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left + seg_w * i, top, seg_w + Emu(5), height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()

def add_rect(slide, left, top, width, height, color, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=WHITE, line_spacing=1.4, font_name="Segoe UI"):
    """lines = list of (text, bold, color) tuples"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.4)
    return txBox

def add_icon_circle(slide, left, top, size, color, icon_text=""):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if icon_text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = icon_text
        p.font.size = Pt(int(Emu(size) / 914400 * 14))
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
    return shape

def add_card(slide, left, top, width, height, title, items, icon_color, icon_text=""):
    """Card premium com ícone, título e bullets"""
    card = add_rect(slide, left, top, width, height, CARD_BG, radius=0.05)

    # Linha de destaque no topo
    accent_bar = add_rect(slide, left, top, width, Inches(0.06), icon_color)

    # Ícone
    add_icon_circle(slide, left + Inches(0.4), top + Inches(0.35), Inches(0.6), icon_color, icon_text)

    # Título
    add_text(slide, left + Inches(1.15), top + Inches(0.3), width - Inches(1.5), Inches(0.5),
             title, font_size=17, color=WHITE, bold=True)

    # Items
    y = top + Inches(1.0)
    for item in items:
        add_text(slide, left + Inches(0.5), y, width - Inches(0.8), Inches(0.35),
                 f"  {item}", font_size=12, color=LIGHT_GRAY)
        y += Inches(0.32)

    return card

def add_metric_card(slide, left, top, width, height, value, label, color):
    """Card de métrica estilo dashboard"""
    card = add_rect(slide, left, top, width, height, CARD_BG, radius=0.05)
    add_text(slide, left, top + Inches(0.25), width, Inches(0.6),
             value, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.85), width, Inches(0.4),
             label, font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    return card

def add_page_number(slide, num, total):
    add_text(slide, Inches(15.0), Inches(8.55), Inches(0.8), Inches(0.3),
             f"{num}/{total}", font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)

def add_bottom_bar(slide):
    add_gradient_bar(slide, Inches(0), Inches(8.85), Inches(16), Inches(0.15),
                     ACCENT_BLUE, ACCENT_CYAN)

TOTAL_SLIDES = 16

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Elemento decorativo - círculo grande semitransparente
c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(10), Inches(10))
c1.fill.solid()
c1.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
c1.line.fill.background()

c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(4), Inches(8), Inches(8))
c2.fill.solid()
c2.fill.fore_color.rgb = RGBColor(0x15, 0x20, 0x3C)
c2.line.fill.background()

# Barra gradiente topo
add_gradient_bar(slide, Inches(0), Inches(0), Inches(16), Inches(0.08), ACCENT_BLUE, ACCENT_CYAN)

# Badge
badge = add_rect(slide, Inches(1.5), Inches(2.2), Inches(2.8), Inches(0.45), ACCENT_BLUE, radius=0.15)
add_text(slide, Inches(1.5), Inches(2.22), Inches(2.8), Inches(0.45),
         "PLATAFORMA ATS", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
         font_name="Segoe UI Semibold")

# Título principal
add_text(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(1.2),
         "ITSector Talent", font_size=54, color=WHITE, bold=True, font_name="Segoe UI Light")
add_text(slide, Inches(1.5), Inches(3.9), Inches(10), Inches(1.2),
         "Inbox ATS", font_size=54, color=ACCENT_CYAN, bold=True, font_name="Segoe UI Light")

# Subtítulo
add_text(slide, Inches(1.5), Inches(5.2), Inches(8), Inches(0.8),
         "Sistema Inteligente de Gestao de Candidaturas e Recrutamento",
         font_size=20, color=LIGHT_GRAY, font_name="Segoe UI")

# Linha separadora
add_rect(slide, Inches(1.5), Inches(6.1), Inches(3), Inches(0.04), ACCENT_BLUE)

add_text(slide, Inches(1.5), Inches(6.4), Inches(6), Inches(0.5),
         "Transformar emails em talento. Automaticamente.",
         font_size=16, color=MED_GRAY, font_name="Segoe UI")

# Três mini-métricas na capa
metrics_y = Inches(7.3)
for i, (val, lbl) in enumerate([("100%", "Local & Seguro"), ("M365", "Integracao Nativa"), ("IA", "Scoring Inteligente")]):
    x = Inches(1.5 + i * 3.2)
    add_text(slide, x, metrics_y, Inches(2), Inches(0.4), val, font_size=22, color=ACCENT_CYAN, bold=True)
    add_text(slide, x, metrics_y + Inches(0.35), Inches(2.5), Inches(0.3), lbl, font_size=12, color=MED_GRAY)

add_bottom_bar(slide)
add_page_number(slide, 1, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — O DESAFIO
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.5), Inches(8), Inches(0.6),
         "O DESAFIO", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(1.0), Inches(12), Inches(0.8),
         "O Recrutamento Tradicional por Email Nao Escala",
         font_size=36, color=WHITE, bold=True, font_name="Segoe UI Light")

# Pain points como cards
pains = [
    ("Emails Perdidos", "Candidaturas enterradas na caixa\nde entrada sem rastreamento", ACCENT_ROSE, "!"),
    ("CVs Manuais", "Leitura manual de dezenas de CVs\npor dia sem normalizacao", ACCENT_AMBER, "?"),
    ("Sem Ranking", "Impossivel comparar candidatos\nobjectivamente entre vagas", ACCENT_PURPLE, "#"),
    ("Tempo Perdido", "Horas gastas em tarefas repetitivas\nque podiam ser automatizadas", ACCENT_ROSE, "T"),
]

for i, (title, desc, color, icon) in enumerate(pains):
    x = Inches(1 + i * 3.6)
    y = Inches(2.5)

    card = add_rect(slide, x, y, Inches(3.2), Inches(3.5), CARD_BG, radius=0.05)
    add_rect(slide, x, y, Inches(3.2), Inches(0.06), color)
    add_icon_circle(slide, x + Inches(1.2), y + Inches(0.4), Inches(0.7), color, icon)
    add_text(slide, x + Inches(0.3), y + Inches(1.3), Inches(2.6), Inches(0.5),
             title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), y + Inches(1.9), Inches(2.6), Inches(1.2),
             desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Estatística impactante
add_rect(slide, Inches(1), Inches(6.8), Inches(14), Inches(1.2), DARK_BG2, radius=0.03)
add_text(slide, Inches(1.5), Inches(7.0), Inches(2), Inches(0.8),
         "67%", font_size=42, color=ACCENT_CYAN, bold=True)
add_text(slide, Inches(3.5), Inches(7.05), Inches(11), Inches(0.8),
         "dos recrutadores afirmam perder candidatos qualificados devido a processos manuais desorganizados",
         font_size=16, color=LIGHT_GRAY)

add_bottom_bar(slide)
add_page_number(slide, 2, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — A SOLUCAO
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.5), Inches(8), Inches(0.6),
         "A SOLUCAO", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(1.0), Inches(13), Inches(0.8),
         "De Email a Contratacao - Tudo Automatizado",
         font_size=36, color=WHITE, bold=True, font_name="Segoe UI Light")

# Workflow horizontal com setas
steps = [
    ("1", "Email Recebido", "Sincronizacao\nautomatica M365", ACCENT_BLUE),
    ("2", "CV Processado", "Parsing PDF/DOCX\ne extracao de dados", ACCENT_CYAN),
    ("3", "Candidato Criado", "Perfil normalizado\ncom competencias", ACCENT_GREEN),
    ("4", "Score Calculado", "Avaliacao heuristica\npor vaga", ACCENT_PURPLE),
    ("5", "Decisao Tomada", "Avancar, Rever\nou Rejeitar", ACCENT_AMBER),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 3.05)
    y = Inches(2.5)

    # Card
    card = add_rect(slide, x, y, Inches(2.7), Inches(3.0), CARD_BG, radius=0.05)
    add_rect(slide, x, y, Inches(2.7), Inches(0.06), color)

    # Número circular
    add_icon_circle(slide, x + Inches(0.95), y + Inches(0.3), Inches(0.7), color, num)

    # Título e descrição
    add_text(slide, x + Inches(0.15), y + Inches(1.2), Inches(2.4), Inches(0.45),
             title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.15), y + Inches(1.7), Inches(2.4), Inches(0.8),
             desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Seta entre cards
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(2.7), y + Inches(1.2), Inches(0.35), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_CYAN
        arrow.line.fill.background()

# Destaque inferior
add_rect(slide, Inches(1), Inches(6.3), Inches(14), Inches(1.8), DARK_BG2, radius=0.03)
add_text(slide, Inches(1.5), Inches(6.5), Inches(13), Inches(0.5),
         "Resultado: Zero Trabalho Manual na Triagem Inicial",
         font_size=22, color=ACCENT_CYAN, bold=True)
add_text(slide, Inches(1.5), Inches(7.1), Inches(13), Inches(0.8),
         "O sistema processa automaticamente cada email, extrai informacao dos CVs, cria perfis de candidatos,\n"
         "calcula pontuacoes de adequacao e recomenda accoes - tudo em segundos.",
         font_size=14, color=LIGHT_GRAY)

add_bottom_bar(slide)
add_page_number(slide, 3, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — DASHBOARD OVERVIEW (Mock Screenshot)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "DASHBOARD", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(10), Inches(0.6),
         "Visao Geral em Tempo Real", font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# Simulated dashboard frame
frame = add_rect(slide, Inches(0.5), Inches(1.5), Inches(15), Inches(7.0), RGBColor(0x0B, 0x11, 0x1F), radius=0.02)

# Sidebar simulada
sidebar = add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.5), Inches(7.0), RGBColor(0x0F, 0x17, 0x2A))

sidebar_items = [
    ("Dashboard", True), ("Inbox", False), ("Candidatos", False),
    ("Vagas", False), ("Talent Pool", False), ("Templates", False),
    ("Relatorios", False), ("Definicoes", False)
]
for i, (name, active) in enumerate(sidebar_items):
    y = Inches(2.0 + i * 0.45)
    if active:
        add_rect(slide, Inches(0.5), y, Inches(2.5), Inches(0.4), ACCENT_BLUE, radius=0.08)
    add_text(slide, Inches(0.9), y + Inches(0.02), Inches(2), Inches(0.35),
             name, font_size=11, color=WHITE if active else MED_GRAY)

# Logo no sidebar
add_text(slide, Inches(0.6), Inches(1.6), Inches(2.2), Inches(0.4),
         "ITSector ATS", font_size=13, color=ACCENT_CYAN, bold=True)

# Métricas principais
metric_data = [
    ("247", "Total Candidatos", ACCENT_BLUE),
    ("12", "Vagas Abertas", ACCENT_GREEN),
    ("38", "Em Revisao", ACCENT_AMBER),
    ("89%", "Taxa Processamento", ACCENT_CYAN),
]
for i, (val, label, color) in enumerate(metric_data):
    x = Inches(3.3 + i * 2.85)
    add_metric_card(slide, x, Inches(1.8), Inches(2.6), Inches(1.5), val, label, color)

# Gráfico simulado (barras)
chart_bg = add_rect(slide, Inches(3.3), Inches(3.6), Inches(5.5), Inches(3.2), CARD_BG, radius=0.03)
add_text(slide, Inches(3.6), Inches(3.7), Inches(4), Inches(0.4),
         "Candidaturas por Vaga", font_size=12, color=LIGHT_GRAY, bold=True)

# Barras simuladas
bars = [
    (".NET Dev", 0.85, ACCENT_BLUE),
    ("React Dev", 0.65, ACCENT_CYAN),
    ("QA Engineer", 0.45, ACCENT_GREEN),
    ("DevOps", 0.55, ACCENT_PURPLE),
    ("Power Apps", 0.35, ACCENT_AMBER),
    ("Mobile Dev", 0.40, ACCENT_ROSE),
]
for i, (label, pct, color) in enumerate(bars):
    by = Inches(4.25 + i * 0.42)
    add_text(slide, Inches(3.6), by, Inches(1.2), Inches(0.3), label, font_size=9, color=MED_GRAY)
    bar_w = Inches(3.5 * pct)
    add_rect(slide, Inches(5.0), by + Inches(0.04), bar_w, Inches(0.22), color, radius=0.2)

# Lista de actividade recente
activity_bg = add_rect(slide, Inches(9.1), Inches(3.6), Inches(6.1), Inches(3.2), CARD_BG, radius=0.03)
add_text(slide, Inches(9.4), Inches(3.7), Inches(4), Inches(0.4),
         "Actividade Recente", font_size=12, color=LIGHT_GRAY, bold=True)

activities = [
    "Novo candidato: Maria Silva (.NET Senior)",
    "CV processado: Joao Santos (React Dev)",
    "Score calculado: Ana Costa - 87/100",
    "Entrevista agendada: Pedro Almeida",
    "Email enviado: Template Rejeicao",
    "Vaga actualizada: QA Engineer -> ABERTA",
]
for i, act in enumerate(activities):
    add_rect(slide, Inches(9.3), Inches(4.2 + i * 0.42), Inches(0.12), Inches(0.12),
             [ACCENT_BLUE, ACCENT_GREEN, ACCENT_CYAN, ACCENT_PURPLE, ACCENT_AMBER, ACCENT_ROSE][i])
    add_text(slide, Inches(9.6), Inches(4.15 + i * 0.42), Inches(5.3), Inches(0.3),
             act, font_size=10, color=LIGHT_GRAY)

# Pipeline mini
pipe_bg = add_rect(slide, Inches(3.3), Inches(7.0), Inches(11.9), Inches(1.2), CARD_BG, radius=0.03)
add_text(slide, Inches(3.6), Inches(7.05), Inches(4), Inches(0.35),
         "Pipeline de Recrutamento", font_size=12, color=LIGHT_GRAY, bold=True)

stages = [("Novo", "42", ACCENT_BLUE), ("Revisao", "28", ACCENT_AMBER), ("Shortlist", "15", ACCENT_CYAN),
          ("Entrevista", "8", ACCENT_PURPLE), ("Oferta", "3", ACCENT_GREEN), ("Contratado", "5", ACCENT_GREEN)]
for i, (name, count, color) in enumerate(stages):
    sx = Inches(3.5 + i * 1.95)
    add_rect(slide, sx, Inches(7.45), Inches(1.7), Inches(0.6), color, radius=0.05)
    add_text(slide, sx, Inches(7.45), Inches(1.7), Inches(0.3),
             count, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, sx, Inches(7.73), Inches(1.7), Inches(0.25),
             name, font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_page_number(slide, 4, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — INBOX & SYNC (Mock)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "INBOX INTELIGENTE", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Sincronizacao Automatica com Microsoft 365",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# Frame da app
frame = add_rect(slide, Inches(0.5), Inches(1.5), Inches(15), Inches(6.8), RGBColor(0x0B, 0x11, 0x1F), radius=0.02)

# Sidebar mini
sidebar = add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.5), Inches(6.8), RGBColor(0x0F, 0x17, 0x2A))
add_text(slide, Inches(0.6), Inches(1.6), Inches(2.2), Inches(0.4),
         "ITSector ATS", font_size=13, color=ACCENT_CYAN, bold=True)

sidebar_items2 = [
    ("Dashboard", False), ("Inbox", True), ("Candidatos", False),
    ("Vagas", False), ("Talent Pool", False)
]
for i, (name, active) in enumerate(sidebar_items2):
    y = Inches(2.2 + i * 0.45)
    if active:
        add_rect(slide, Inches(0.5), y, Inches(2.5), Inches(0.4), ACCENT_BLUE, radius=0.08)
    add_text(slide, Inches(0.9), y + Inches(0.02), Inches(2), Inches(0.35),
             name, font_size=11, color=WHITE if active else MED_GRAY)

# Header da inbox
add_text(slide, Inches(3.3), Inches(1.7), Inches(5), Inches(0.4),
         "Inbox - careers@itsector.pt", font_size=16, color=WHITE, bold=True)

# Botão sync
sync_btn = add_rect(slide, Inches(12.5), Inches(1.65), Inches(2.5), Inches(0.5), ACCENT_BLUE, radius=0.1)
add_text(slide, Inches(12.5), Inches(1.68), Inches(2.5), Inches(0.45),
         "Sincronizar Agora", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Status bar
add_rect(slide, Inches(3.3), Inches(2.3), Inches(12), Inches(0.45), DARK_BG2, radius=0.05)
add_text(slide, Inches(3.6), Inches(2.32), Inches(10), Inches(0.35),
         "Ultima sincronizacao: ha 5 minutos  |  23 emails processados  |  8 novos candidatos criados",
         font_size=10, color=ACCENT_GREEN)

# Lista de emails
emails = [
    ("Maria Silva", "Candidatura - .NET Developer Senior", "14:32", "CV Processado", ACCENT_GREEN, True),
    ("Joao Santos", "Candidatura React Developer - Portfolio", "13:45", "Em Analise", ACCENT_AMBER, True),
    ("Ana Costa", "Re: Oportunidade QA Engineer", "12:20", "Novo", ACCENT_BLUE, True),
    ("Pedro Almeida", "Candidatura Espontanea - DevOps", "11:05", "CV Processado", ACCENT_GREEN, True),
    ("Newsletter SAPO", "Tech News Semanal", "10:30", "Ignorado", MED_GRAY, False),
    ("Sofia Martins", "Interesse vaga Power Apps", "09:15", "Novo", ACCENT_BLUE, True),
    ("Carlos Ferreira", "CV em anexo - Mobile Developer", "08:40", "CV Processado", ACCENT_GREEN, True),
]

for i, (name, subject, time, status, status_color, is_candidate) in enumerate(emails):
    ey = Inches(2.95 + i * 0.65)
    row_color = CARD_BG if i % 2 == 0 else RGBColor(0x17, 0x22, 0x34)
    add_rect(slide, Inches(3.3), ey, Inches(12), Inches(0.6), row_color, radius=0.02)

    # Indicador
    if is_candidate:
        add_rect(slide, Inches(3.3), ey, Inches(0.06), Inches(0.6), ACCENT_GREEN)

    add_text(slide, Inches(3.6), ey + Inches(0.05), Inches(2), Inches(0.25),
             name, font_size=11, color=WHITE, bold=True)
    add_text(slide, Inches(5.6), ey + Inches(0.05), Inches(5.5), Inches(0.25),
             subject, font_size=11, color=LIGHT_GRAY)
    add_text(slide, Inches(12.5), ey + Inches(0.05), Inches(0.8), Inches(0.25),
             time, font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)

    # Badge status
    badge_w = Inches(1.3)
    badge_bg = add_rect(slide, Inches(13.5), ey + Inches(0.1), badge_w, Inches(0.32),
                        darken(status_color), radius=0.15)
    add_text(slide, Inches(13.5), ey + Inches(0.12), badge_w, Inches(0.3),
             status, font_size=9, color=status_color, bold=True, alignment=PP_ALIGN.CENTER)

# Legenda
add_text(slide, Inches(3.5), Inches(7.75), Inches(3), Inches(0.3),
         "Verde = Candidatura Detectada", font_size=10, color=ACCENT_GREEN)
add_text(slide, Inches(7), Inches(7.75), Inches(3), Inches(0.3),
         "Cinza = Email Ignorado (nao-candidatura)", font_size=10, color=MED_GRAY)

add_bottom_bar(slide)
add_page_number(slide, 5, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — GESTAO DE CANDIDATOS (Mock)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "GESTAO DE CANDIDATOS", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Perfis Completos Extraidos Automaticamente",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# Candidate detail mock
frame = add_rect(slide, Inches(0.5), Inches(1.5), Inches(15), Inches(6.8), RGBColor(0x0B, 0x11, 0x1F), radius=0.02)

# Perfil header
add_rect(slide, Inches(0.8), Inches(1.7), Inches(14.4), Inches(1.8), CARD_BG, radius=0.03)

# Avatar circle
add_icon_circle(slide, Inches(1.2), Inches(1.9), Inches(1.2), ACCENT_BLUE, "MS")

add_text(slide, Inches(2.7), Inches(1.9), Inches(5), Inches(0.45),
         "Maria Silva", font_size=22, color=WHITE, bold=True)
add_text(slide, Inches(2.7), Inches(2.35), Inches(5), Inches(0.35),
         "Senior .NET Developer  |  Lisboa  |  8 anos experiencia", font_size=12, color=LIGHT_GRAY)

# Score badge
score_bg = add_rect(slide, Inches(12.5), Inches(1.9), Inches(2.2), Inches(1.3), ACCENT_GREEN, radius=0.05)
add_text(slide, Inches(12.5), Inches(1.95), Inches(2.2), Inches(0.7),
         "87", font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(12.5), Inches(2.7), Inches(2.2), Inches(0.3),
         "AVANCAR", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Status badges
statuses_x = Inches(8)
for j, (lbl, clr) in enumerate([("Shortlisted", ACCENT_CYAN), (".NET Developer", ACCENT_BLUE), ("Entrevista Agendada", ACCENT_PURPLE)]):
    bx = statuses_x + Inches(j * 2.0)
    add_rect(slide, bx, Inches(2.9), Inches(1.8), Inches(0.3),
             darken(clr), radius=0.15)
    add_text(slide, bx, Inches(2.91), Inches(1.8), Inches(0.28),
             lbl, font_size=9, color=clr, bold=True, alignment=PP_ALIGN.CENTER)

# Secção de competências
add_rect(slide, Inches(0.8), Inches(3.8), Inches(7), Inches(2.3), CARD_BG, radius=0.03)
add_text(slide, Inches(1.1), Inches(3.9), Inches(5), Inches(0.35),
         "Competencias Extraidas do CV", font_size=13, color=WHITE, bold=True)

skills = ["C#/.NET", "Azure", "SQL Server", "REST APIs", "Microservices", "Docker", "CI/CD", "Agile"]
for i, skill in enumerate(skills):
    row = i // 4
    col = i % 4
    sx = Inches(1.1 + col * 1.7)
    sy = Inches(4.4 + row * 0.5)
    add_rect(slide, sx, sy, Inches(1.5), Inches(0.38),
             darken(ACCENT_BLUE), radius=0.15)
    add_text(slide, sx, sy + Inches(0.02), Inches(1.5), Inches(0.35),
             skill, font_size=10, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Score breakdown
add_rect(slide, Inches(8.1), Inches(3.8), Inches(7), Inches(2.3), CARD_BG, radius=0.03)
add_text(slide, Inches(8.4), Inches(3.9), Inches(5), Inches(0.35),
         "Decomposicao do Score", font_size=13, color=WHITE, bold=True)

scores = [
    ("Skills Obrigatorias", "28/30", 0.93, ACCENT_GREEN),
    ("Skills Opcionais", "8/10", 0.80, ACCENT_CYAN),
    ("Experiencia", "13/15", 0.87, ACCENT_BLUE),
    ("Idiomas", "10/10", 1.0, ACCENT_GREEN),
    ("Localizacao", "5/5", 1.0, ACCENT_GREEN),
]
for i, (label, val, pct, color) in enumerate(scores):
    sy = Inches(4.4 + i * 0.38)
    add_text(slide, Inches(8.4), sy, Inches(2.5), Inches(0.3), label, font_size=10, color=LIGHT_GRAY)
    bar_bg = add_rect(slide, Inches(11.2), sy + Inches(0.05), Inches(2.5), Inches(0.18), DARK_BG2, radius=0.3)
    bar_fill = add_rect(slide, Inches(11.2), sy + Inches(0.05), Inches(2.5 * pct), Inches(0.18), color, radius=0.3)
    add_text(slide, Inches(13.8), sy, Inches(0.8), Inches(0.3), val, font_size=10, color=color, alignment=PP_ALIGN.RIGHT)

# Secção inferior - info e acções
add_rect(slide, Inches(0.8), Inches(6.3), Inches(7), Inches(1.7), CARD_BG, radius=0.03)
add_text(slide, Inches(1.1), Inches(6.4), Inches(5), Inches(0.35),
         "Informacao Pessoal", font_size=13, color=WHITE, bold=True)
info_items = ["Email: maria.silva@gmail.com", "Telefone: +351 912 345 678",
              "LinkedIn: /in/maria-silva-dev", "Educacao: MSc Eng. Informatica - IST"]
for i, item in enumerate(info_items):
    add_text(slide, Inches(1.1), Inches(6.85 + i * 0.3), Inches(6), Inches(0.28),
             item, font_size=10, color=LIGHT_GRAY)

# Acções
add_rect(slide, Inches(8.1), Inches(6.3), Inches(7), Inches(1.7), CARD_BG, radius=0.03)
add_text(slide, Inches(8.4), Inches(6.4), Inches(5), Inches(0.35),
         "Accoes Rapidas", font_size=13, color=WHITE, bold=True)

actions = [("Enviar Email", ACCENT_BLUE), ("Agendar Entrevista", ACCENT_PURPLE),
           ("Adicionar ao Talent Pool", ACCENT_GREEN), ("Exportar Perfil", ACCENT_CYAN)]
for i, (action, color) in enumerate(actions):
    row = i // 2
    col = i % 2
    ax = Inches(8.4 + col * 3.3)
    ay = Inches(6.85 + row * 0.5)
    add_rect(slide, ax, ay, Inches(2.8), Inches(0.4), color, radius=0.08)
    add_text(slide, ax, ay + Inches(0.03), Inches(2.8), Inches(0.35),
             action, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_page_number(slide, 6, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — CV PARSING
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "PROCESSAMENTO DE CVs", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Extracao Inteligente de Curriculos",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# Left side - CV Document mock
add_rect(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(6.5), RGBColor(0xF8, 0xFA, 0xFC), radius=0.03)
add_text(slide, Inches(1.2), Inches(1.9), Inches(4), Inches(0.4),
         "curriculum_vitae.pdf", font_size=14, color=RGBColor(0x1E, 0x29, 0x3B), bold=True)
add_rect(slide, Inches(1.2), Inches(2.35), Inches(4.2), Inches(0.04), RGBColor(0xE2, 0xE8, 0xF0))

cv_lines = [
    "MARIA SILVA",
    "Senior .NET Developer",
    "",
    "EXPERIENCIA PROFISSIONAL",
    "- ITSector (2019-presente) - .NET Developer",
    "- BNP Paribas (2016-2019) - Software Engineer",
    "",
    "COMPETENCIAS TECNICAS",
    "C#, .NET Core, Azure, SQL Server, Docker",
    "REST APIs, Microservices, CI/CD, Git",
    "",
    "IDIOMAS",
    "Portugues (Nativo), Ingles (Fluente)",
    "",
    "FORMACAO",
    "MSc Engenharia Informatica - IST (2016)",
]
for i, line in enumerate(cv_lines):
    is_header = line.isupper() and line != ""
    add_text(slide, Inches(1.3), Inches(2.5 + i * 0.3), Inches(4), Inches(0.28),
             line, font_size=9 if not is_header else 10,
             color=RGBColor(0x0F, 0x17, 0x2A) if is_header else RGBColor(0x47, 0x55, 0x69),
             bold=is_header)

# Arrow
arrow = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_ARROW, Inches(6.0), Inches(4.2), Inches(1.2), Inches(0.6)
)
arrow.fill.solid()
arrow.fill.fore_color.rgb = ACCENT_CYAN
arrow.line.fill.background()
add_text(slide, Inches(6.0), Inches(4.9), Inches(1.2), Inches(0.3),
         "Parsing", font_size=10, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Right side - Extracted data
add_rect(slide, Inches(7.5), Inches(1.7), Inches(7.5), Inches(6.5), CARD_BG, radius=0.03)
add_text(slide, Inches(7.8), Inches(1.9), Inches(5), Inches(0.4),
         "Dados Extraidos Automaticamente", font_size=14, color=ACCENT_GREEN, bold=True)

extracted = [
    ("Nome", "Maria Silva", ACCENT_CYAN),
    ("Titulo Actual", "Senior .NET Developer", ACCENT_CYAN),
    ("Localizacao", "Lisboa, Portugal", ACCENT_CYAN),
    ("Anos Experiencia", "8 anos", ACCENT_CYAN),
    ("Skills", '["C#", ".NET Core", "Azure", "SQL Server", "Docker"]', ACCENT_GREEN),
    ("Idiomas", '["Portugues (Nativo)", "Ingles (Fluente)"]', ACCENT_GREEN),
    ("Educacao", '["MSc Eng. Informatica - IST 2016"]', ACCENT_GREEN),
    ("Experiencia", '["ITSector 2019-2025", "BNP Paribas 2016-2019"]', ACCENT_GREEN),
    ("LinkedIn", "linkedin.com/in/maria-silva-dev", ACCENT_BLUE),
    ("Email", "maria.silva@gmail.com", ACCENT_BLUE),
    ("Confianca Parser", "94%", ACCENT_GREEN),
]
for i, (key, val, color) in enumerate(extracted):
    ey = Inches(2.5 + i * 0.47)
    add_text(slide, Inches(7.9), ey, Inches(2.2), Inches(0.3),
             key, font_size=10, color=MED_GRAY, bold=True)
    add_text(slide, Inches(10.2), ey, Inches(4.5), Inches(0.3),
             val, font_size=10, color=color)

# Formatos suportados
add_rect(slide, Inches(7.8), Inches(7.2), Inches(6.8), Inches(0.7), DARK_BG2, radius=0.05)
formats = [("PDF", ACCENT_BLUE), ("DOCX", ACCENT_GREEN), ("DOC", ACCENT_CYAN), ("OCR", ACCENT_PURPLE)]
for i, (fmt, clr) in enumerate(formats):
    fx = Inches(8.2 + i * 1.6)
    add_rect(slide, fx, Inches(7.35), Inches(1.2), Inches(0.4), clr, radius=0.1)
    add_text(slide, fx, Inches(7.37), Inches(1.2), Inches(0.38),
             fmt, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_page_number(slide, 7, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — SCORING ENGINE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "MOTOR DE SCORING", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Avaliacao Objectiva e Configuravel",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# Score wheel (circle representations)
center_x, center_y = Inches(3.5), Inches(4.5)
# Outer ring
add_icon_circle(slide, center_x - Inches(1.5), center_y - Inches(1.5), Inches(3), ACCENT_BLUE, "")
add_icon_circle(slide, center_x - Inches(1.2), center_y - Inches(1.2), Inches(2.4), DARK_BG, "")
add_text(slide, center_x - Inches(0.8), center_y - Inches(0.5), Inches(1.6), Inches(0.6),
         "100", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, center_x - Inches(0.5), center_y + Inches(0.3), Inches(1), Inches(0.3),
         "pontos", font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Score components
components = [
    ("Skills Obrigatorias", "30 pts", "Match com requisitos da vaga", ACCENT_BLUE, 0.30),
    ("Anos Experiencia", "15 pts", "Adequacao ao nivel exigido", ACCENT_CYAN, 0.15),
    ("Skills Opcionais", "10 pts", "Competencias diferenciadoras", ACCENT_GREEN, 0.10),
    ("Idiomas", "10 pts", "Correspondencia linguistica", ACCENT_PURPLE, 0.10),
    ("Titulo/Senioridade", "10 pts", "Relevancia funcional", ACCENT_AMBER, 0.10),
    ("Dominio", "10 pts", "Experiencia no sector", ACCENT_ROSE, 0.10),
    ("Localizacao", "5 pts", "Proximidade geografica", ACCENT_CYAN, 0.05),
    ("Qualidade CV", "5 pts", "Completude do curriculo", MED_GRAY, 0.05),
    ("Factor Recrutador", "5 pts", "Ajuste manual do recruiter", ACCENT_BLUE, 0.05),
]

for i, (name, pts, desc, color, pct) in enumerate(components):
    cy = Inches(1.8 + i * 0.65)
    cx = Inches(7)

    add_rect(slide, cx, cy, Inches(0.4), Inches(0.4), color, radius=0.15)
    add_text(slide, cx + Inches(0.55), cy, Inches(3), Inches(0.25),
             name, font_size=12, color=WHITE, bold=True)
    add_text(slide, cx + Inches(0.55), cy + Inches(0.22), Inches(4), Inches(0.2),
             desc, font_size=9, color=MED_GRAY)

    # Barra proporcional
    add_rect(slide, Inches(11.5), cy + Inches(0.08), Inches(3), Inches(0.14), DARK_BG2, radius=0.3)
    add_rect(slide, Inches(11.5), cy + Inches(0.08), Inches(3 * pct / 0.30), Inches(0.14), color, radius=0.3)
    add_text(slide, Inches(14.6), cy, Inches(0.8), Inches(0.3),
             pts, font_size=11, color=color, bold=True, alignment=PP_ALIGN.RIGHT)

# Thresholds
add_rect(slide, Inches(0.8), Inches(7.3), Inches(14.4), Inches(0.9), DARK_BG2, radius=0.03)
thresholds = [
    ("75+", "AVANCAR", "Auto-avanco para entrevista", ACCENT_GREEN),
    ("45-74", "REVISAO MANUAL", "Requer decisao do recrutador", ACCENT_AMBER),
    ("<45", "REJEITAR", "Abaixo do minimo para a vaga", ACCENT_ROSE),
]
for i, (score, action, desc, color) in enumerate(thresholds):
    tx = Inches(1.0 + i * 4.8)
    add_rect(slide, tx, Inches(7.4), Inches(0.8), Inches(0.6), color, radius=0.1)
    add_text(slide, tx, Inches(7.43), Inches(0.8), Inches(0.55),
             score, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, tx + Inches(1.0), Inches(7.4), Inches(2.5), Inches(0.3),
             action, font_size=12, color=color, bold=True)
    add_text(slide, tx + Inches(1.0), Inches(7.68), Inches(3.5), Inches(0.25),
             desc, font_size=10, color=MED_GRAY)

add_bottom_bar(slide)
add_page_number(slide, 8, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — GESTAO DE VAGAS (Mock)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "GESTAO DE VAGAS", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Vagas com Ranking Automatico de Candidatos",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# Jobs grid
jobs = [
    (".NET Developer", "Tecnologia", "Lisboa", "ABERTA", 18, ACCENT_GREEN),
    ("Senior .NET Developer", "Tecnologia", "Porto", "ABERTA", 12, ACCENT_GREEN),
    ("React Developer", "Frontend", "Remoto", "ABERTA", 15, ACCENT_GREEN),
    ("Power Apps Developer", "Low Code", "Lisboa", "ABERTA", 6, ACCENT_GREEN),
    ("Business Analyst", "Consultoria", "Lisboa", "ABERTA", 9, ACCENT_GREEN),
    ("QA Engineer", "Qualidade", "Porto", "PAUSADA", 7, ACCENT_AMBER),
    ("DevOps Engineer", "Infraestrutura", "Remoto", "ABERTA", 11, ACCENT_GREEN),
    ("Mobile Developer", "Mobile", "Lisboa", "RASCUNHO", 3, MED_GRAY),
]

for i, (title, dept, loc, status, applicants, status_color) in enumerate(jobs):
    row = i // 4
    col = i % 4
    cx = Inches(0.5 + col * 3.85)
    cy = Inches(1.7 + row * 3.3)

    card = add_rect(slide, cx, cy, Inches(3.6), Inches(3.0), CARD_BG, radius=0.04)
    add_rect(slide, cx, cy, Inches(3.6), Inches(0.05), status_color)

    add_text(slide, cx + Inches(0.3), cy + Inches(0.2), Inches(3), Inches(0.35),
             title, font_size=14, color=WHITE, bold=True)
    add_text(slide, cx + Inches(0.3), cy + Inches(0.55), Inches(3), Inches(0.25),
             f"{dept}  |  {loc}", font_size=10, color=MED_GRAY)

    # Status badge
    add_rect(slide, cx + Inches(0.3), cy + Inches(0.9), Inches(1.0), Inches(0.28),
             darken(status_color), radius=0.15)
    add_text(slide, cx + Inches(0.3), cy + Inches(0.91), Inches(1.0), Inches(0.27),
             status, font_size=8, color=status_color, bold=True, alignment=PP_ALIGN.CENTER)

    # Applicants count
    add_text(slide, cx + Inches(0.3), cy + Inches(1.4), Inches(1), Inches(0.5),
             str(applicants), font_size=28, color=ACCENT_CYAN, bold=True)
    add_text(slide, cx + Inches(1.3), cy + Inches(1.5), Inches(2), Inches(0.3),
             "candidatos", font_size=11, color=MED_GRAY)

    # Mini ranking bars
    for j in range(min(3, applicants)):
        bar_pct = [0.95, 0.78, 0.62][j]
        by = cy + Inches(2.0 + j * 0.28)
        add_text(slide, cx + Inches(0.3), by, Inches(0.3), Inches(0.2),
                 f"#{j+1}", font_size=8, color=MED_GRAY)
        add_rect(slide, cx + Inches(0.65), by + Inches(0.03), Inches(2.2), Inches(0.14), DARK_BG2, radius=0.3)
        add_rect(slide, cx + Inches(0.65), by + Inches(0.03), Inches(2.2 * bar_pct), Inches(0.14),
                 [ACCENT_GREEN, ACCENT_CYAN, ACCENT_AMBER][j], radius=0.3)

add_bottom_bar(slide)
add_page_number(slide, 9, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — EMAIL & COMUNICACAO (Mock)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "COMUNICACAO INTEGRADA", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Templates, Emails e Agendamento de Entrevistas",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# Left - Templates
add_rect(slide, Inches(0.5), Inches(1.6), Inches(7.2), Inches(6.5), CARD_BG, radius=0.03)
add_text(slide, Inches(0.9), Inches(1.8), Inches(5), Inches(0.4),
         "Templates de Email", font_size=16, color=WHITE, bold=True)

templates = [
    ("Convite para Entrevista", "Agenda automaticamente via M365 Calendar", ACCENT_PURPLE),
    ("Rejeicao Profissional", "Resposta respeitosa com feedback", ACCENT_ROSE),
    ("Proposta/Oferta", "Template com detalhes da posicao", ACCENT_GREEN),
    ("Pedido de Informacao", "Solicitar dados adicionais ao candidato", ACCENT_AMBER),
    ("Boas-vindas Talent Pool", "Informar inclusao no talent pool", ACCENT_CYAN),
]

for i, (name, desc, color) in enumerate(templates):
    ty = Inches(2.5 + i * 0.95)
    add_rect(slide, Inches(0.8), ty, Inches(6.6), Inches(0.8), DARK_BG2, radius=0.03)
    add_rect(slide, Inches(0.8), ty, Inches(0.08), Inches(0.8), color)
    add_text(slide, Inches(1.2), ty + Inches(0.08), Inches(5), Inches(0.3),
             name, font_size=12, color=WHITE, bold=True)
    add_text(slide, Inches(1.2), ty + Inches(0.38), Inches(5.5), Inches(0.3),
             desc, font_size=10, color=MED_GRAY)
    # Use button
    add_rect(slide, Inches(6.2), ty + Inches(0.2), Inches(1), Inches(0.35), color, radius=0.1)
    add_text(slide, Inches(6.2), ty + Inches(0.22), Inches(1), Inches(0.33),
             "Usar", font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Right - Email composer mock
add_rect(slide, Inches(8.2), Inches(1.6), Inches(7.2), Inches(6.5), CARD_BG, radius=0.03)
add_text(slide, Inches(8.5), Inches(1.8), Inches(5), Inches(0.4),
         "Compor Email", font_size=16, color=WHITE, bold=True)

# Form fields
fields = [
    ("Para:", "maria.silva@gmail.com"),
    ("Assunto:", "Convite para Entrevista - .NET Developer"),
]
for i, (label, value) in enumerate(fields):
    fy = Inches(2.5 + i * 0.6)
    add_text(slide, Inches(8.5), fy, Inches(1), Inches(0.3), label, font_size=11, color=MED_GRAY)
    add_rect(slide, Inches(9.5), fy, Inches(5.5), Inches(0.4), DARK_BG2, radius=0.05)
    add_text(slide, Inches(9.7), fy + Inches(0.03), Inches(5), Inches(0.3),
             value, font_size=11, color=LIGHT_GRAY)

# Body
add_rect(slide, Inches(8.5), Inches(3.9), Inches(6.6), Inches(3.0), DARK_BG2, radius=0.05)
body_lines = [
    "Cara Maria Silva,",
    "",
    "Temos o prazer de a convidar para uma",
    "entrevista relativa a posicao de .NET Developer",
    "na ITSector.",
    "",
    "Data: 28 de Marco de 2026, 14:30",
    "Local: Escritorio ITSector Lisboa",
    "",
    "Aguardamos a sua confirmacao.",
    "Com os melhores cumprimentos,",
    "Equipa de Recrutamento ITSector",
]
for i, line in enumerate(body_lines):
    add_text(slide, Inches(8.7), Inches(4.0 + i * 0.24), Inches(6), Inches(0.22),
             line, font_size=10, color=LIGHT_GRAY)

# Send button
add_rect(slide, Inches(13.2), Inches(7.2), Inches(1.8), Inches(0.5), ACCENT_BLUE, radius=0.1)
add_text(slide, Inches(13.2), Inches(7.23), Inches(1.8), Inches(0.45),
         "Enviar via M365", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Calendar button
add_rect(slide, Inches(11.2), Inches(7.2), Inches(1.8), Inches(0.5), ACCENT_PURPLE, radius=0.1)
add_text(slide, Inches(11.2), Inches(7.23), Inches(1.8), Inches(0.45),
         "Agendar M365", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_page_number(slide, 10, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 11 — TALENT POOL
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "TALENT POOL", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Base de Talento Reutilizavel para Futuras Oportunidades",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

add_text(slide, Inches(1), Inches(1.5), Inches(13), Inches(0.5),
         "Candidatos qualificados que nao encaixam na vaga actual sao preservados para futuras oportunidades.",
         font_size=14, color=LIGHT_GRAY)

# Pool cards
pool_candidates = [
    ("Ana Rodrigues", "React Developer Junior", "Score: 58 (Near-Fit QA)", ["React", "TypeScript", "Node.js"],
     "Sugestao: React Developer (quando abrir junior)", ACCENT_CYAN),
    ("Bruno Mendes", "DevOps Engineer", "Score: 62 (Near-Fit)", ["AWS", "Terraform", "Kubernetes"],
     "Sugestao: Proxima vaga Cloud/Infra", ACCENT_PURPLE),
    ("Carla Dias", "Business Analyst", "Score: 71 (Revisao Manual)", ["BPMN", "SQL", "Power BI"],
     "Sugestao: Rever quando abrir BA Senior", ACCENT_AMBER),
    ("Daniel Sousa", "Full-Stack Developer", "Candidatura Espontanea", [".NET", "Angular", "Azure"],
     "Sugestao: Match com .NET Developer ou React Dev", ACCENT_GREEN),
]

for i, (name, role, score_info, skills, suggestion, color) in enumerate(pool_candidates):
    cx = Inches(0.5 + (i % 2) * 7.7)
    cy = Inches(2.2 + (i // 2) * 2.8)

    card = add_rect(slide, cx, cy, Inches(7.2), Inches(2.5), CARD_BG, radius=0.04)
    add_rect(slide, cx, cy, Inches(7.2), Inches(0.05), color)

    # Avatar
    initials = "".join([w[0] for w in name.split()])
    add_icon_circle(slide, cx + Inches(0.3), cy + Inches(0.3), Inches(0.7), color, initials)

    add_text(slide, cx + Inches(1.2), cy + Inches(0.3), Inches(4), Inches(0.3),
             name, font_size=14, color=WHITE, bold=True)
    add_text(slide, cx + Inches(1.2), cy + Inches(0.6), Inches(4), Inches(0.25),
             role, font_size=11, color=MED_GRAY)
    add_text(slide, cx + Inches(1.2), cy + Inches(0.85), Inches(4), Inches(0.25),
             score_info, font_size=10, color=color)

    # Skills
    for j, skill in enumerate(skills):
        sx = cx + Inches(0.3 + j * 1.3)
        add_rect(slide, sx, cy + Inches(1.25), Inches(1.1), Inches(0.3),
                 darken(color), radius=0.15)
        add_text(slide, sx, cy + Inches(1.27), Inches(1.1), Inches(0.28),
                 skill, font_size=9, color=color, alignment=PP_ALIGN.CENTER)

    # Suggestion
    add_rect(slide, cx + Inches(0.3), cy + Inches(1.75), Inches(6.5), Inches(0.5), DARK_BG2, radius=0.04)
    add_text(slide, cx + Inches(0.5), cy + Inches(1.8), Inches(6), Inches(0.4),
             suggestion, font_size=10, color=ACCENT_GREEN)

# Bottom insight
add_rect(slide, Inches(0.5), Inches(7.8), Inches(15), Inches(0.6), DARK_BG2, radius=0.03)
add_text(slide, Inches(1), Inches(7.85), Inches(14), Inches(0.5),
         "O Talent Pool garante que nenhum candidato qualificado se perde. "
         "Sugestoes automaticas de vagas futuras maximizam o retorno do pipeline.",
         font_size=12, color=LIGHT_GRAY)

add_bottom_bar(slide)
add_page_number(slide, 11, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 12 — DUPLICADOS & QUALIDADE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "QUALIDADE DE DADOS", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Deteccao de Duplicados e Exportacao",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# Left - Duplicate detection
add_rect(slide, Inches(0.5), Inches(1.7), Inches(7.2), Inches(5.5), CARD_BG, radius=0.03)
add_text(slide, Inches(0.9), Inches(1.9), Inches(5), Inches(0.4),
         "Deteccao Inteligente de Duplicados", font_size=16, color=WHITE, bold=True)

add_text(slide, Inches(0.9), Inches(2.5), Inches(6.5), Inches(0.6),
         "O sistema detecta automaticamente candidatos duplicados\ncom base em email, nome e telefone.",
         font_size=12, color=LIGHT_GRAY)

# Duplicate example
add_rect(slide, Inches(0.9), Inches(3.3), Inches(6.4), Inches(3.5), DARK_BG2, radius=0.03)
add_text(slide, Inches(1.2), Inches(3.4), Inches(5), Inches(0.35),
         "Duplicado Detectado", font_size=13, color=ACCENT_AMBER, bold=True)

# Record A
add_rect(slide, Inches(1.2), Inches(3.9), Inches(2.8), Inches(2.0), CARD_BG, radius=0.03)
add_text(slide, Inches(1.4), Inches(4.0), Inches(2.5), Inches(0.3),
         "Registo A", font_size=11, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1.4), Inches(4.3), Inches(2.5), Inches(0.2),
         "Joao Santos", font_size=10, color=WHITE)
add_text(slide, Inches(1.4), Inches(4.55), Inches(2.5), Inches(0.2),
         "joao.santos@email.com", font_size=9, color=MED_GRAY)
add_text(slide, Inches(1.4), Inches(4.8), Inches(2.5), Inches(0.2),
         "Via: Candidatura directa", font_size=9, color=MED_GRAY)
add_text(slide, Inches(1.4), Inches(5.1), Inches(2.5), Inches(0.2),
         "Score: 72", font_size=10, color=ACCENT_AMBER)

# Merge arrow
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(4.7), Inches(0.5), Inches(0.3))
arrow.fill.solid()
arrow.fill.fore_color.rgb = ACCENT_GREEN
arrow.line.fill.background()

# Record B
add_rect(slide, Inches(4.9), Inches(3.9), Inches(2.8), Inches(2.0), CARD_BG, radius=0.03)
add_text(slide, Inches(5.1), Inches(4.0), Inches(2.5), Inches(0.3),
         "Registo B", font_size=11, color=ACCENT_PURPLE, bold=True)
add_text(slide, Inches(5.1), Inches(4.3), Inches(2.5), Inches(0.2),
         "J. Santos", font_size=10, color=WHITE)
add_text(slide, Inches(5.1), Inches(4.55), Inches(2.5), Inches(0.2),
         "joao.santos@email.com", font_size=9, color=MED_GRAY)
add_text(slide, Inches(5.1), Inches(4.8), Inches(2.5), Inches(0.2),
         "Via: Email espontaneo", font_size=9, color=MED_GRAY)
add_text(slide, Inches(5.1), Inches(5.1), Inches(2.5), Inches(0.2),
         "Score: 65", font_size=10, color=ACCENT_AMBER)

# Merge button
add_rect(slide, Inches(2.5), Inches(6.1), Inches(3), Inches(0.45), ACCENT_GREEN, radius=0.1)
add_text(slide, Inches(2.5), Inches(6.13), Inches(3), Inches(0.4),
         "Fundir Registos", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Right - Exports and Reports
add_rect(slide, Inches(8.2), Inches(1.7), Inches(7.2), Inches(5.5), CARD_BG, radius=0.03)
add_text(slide, Inches(8.5), Inches(1.9), Inches(5), Inches(0.4),
         "Relatorios e Exportacoes", font_size=16, color=WHITE, bold=True)

exports = [
    ("Exportacao Excel", "Lista completa de candidatos com\ntodos os dados extraidos e scores", "XLSX", ACCENT_GREEN),
    ("Relatorio por Vaga", "Ranking de candidatos e metricas\nde cada vaga aberta", "XLSX", ACCENT_BLUE),
    ("Audit Log", "Registo de todas as accoes\nsensitivas no sistema", "LOG", ACCENT_AMBER),
    ("Dashboard Analytics", "Graficos interactivos com KPIs\nde recrutamento em tempo real", "WEB", ACCENT_CYAN),
]

for i, (name, desc, fmt, color) in enumerate(exports):
    ey = Inches(2.5 + i * 1.15)
    add_rect(slide, Inches(8.5), ey, Inches(6.5), Inches(1.0), DARK_BG2, radius=0.03)
    add_rect(slide, Inches(8.5), ey, Inches(0.08), Inches(1.0), color)
    add_text(slide, Inches(8.9), ey + Inches(0.1), Inches(4), Inches(0.3),
             name, font_size=13, color=WHITE, bold=True)
    add_text(slide, Inches(8.9), ey + Inches(0.4), Inches(4.5), Inches(0.5),
             desc, font_size=10, color=MED_GRAY)
    add_rect(slide, Inches(14.0), ey + Inches(0.3), Inches(0.8), Inches(0.35), color, radius=0.1)
    add_text(slide, Inches(14.0), ey + Inches(0.32), Inches(0.8), Inches(0.33),
             fmt, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_page_number(slide, 12, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 13 — INTEGRACOES
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "INTEGRACOES", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Ecossistema Microsoft 365 e Mais",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# Integration cards
integrations = [
    ("Microsoft Graph\nMail", "Sincronizacao bidirecional de\nemails com caixa de correio\ncareers@itsector.pt\n\nMail.Read | Mail.Send\nMail.ReadWrite",
     ACCENT_BLUE, "M"),
    ("Microsoft Graph\nCalendar", "Agendamento de entrevistas\ndirectamente no calendario\ndo Microsoft 365\n\nCalendars.ReadWrite\nConvites automaticos",
     ACCENT_PURPLE, "C"),
    ("Tesseract.js\nOCR", "Reconhecimento optico de\ncaracteres para CVs digitalizados\nem formato imagem\n\nMultiplos idiomas\nFallback automatico",
     ACCENT_CYAN, "O"),
    ("ExcelJS\nExportacao", "Geracao de relatorios Excel\nprofissionais com dados de\ncandidatos e metricas\n\nExportacao local\nFormatos configurados",
     ACCENT_GREEN, "X"),
]

for i, (title, desc, color, icon) in enumerate(integrations):
    cx = Inches(0.5 + i * 3.85)
    cy = Inches(1.8)

    card = add_rect(slide, cx, cy, Inches(3.6), Inches(4.5), CARD_BG, radius=0.04)
    add_rect(slide, cx, cy, Inches(3.6), Inches(0.06), color)

    add_icon_circle(slide, cx + Inches(1.35), cy + Inches(0.4), Inches(0.85), color, icon)

    add_text(slide, cx + Inches(0.3), cy + Inches(1.5), Inches(3), Inches(0.7),
             title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, cx + Inches(0.3), cy + Inches(2.3), Inches(3), Inches(2.0),
             desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# AI section
add_rect(slide, Inches(0.5), Inches(6.7), Inches(15), Inches(1.5), CARD_BG, radius=0.03)
add_text(slide, Inches(1), Inches(6.85), Inches(5), Inches(0.4),
         "Inteligencia Artificial (Opcional)", font_size=16, color=ACCENT_PURPLE, bold=True)
add_text(slide, Inches(1), Inches(7.3), Inches(13), Inches(0.7),
         "Suporte para OpenAI, Azure OpenAI e Anthropic - para resumos de CV, sugestoes de email e classificacao avancada.\n"
         "Totalmente opcional: o sistema funciona com heuristicas locais sem necessidade de API externa.",
         font_size=12, color=LIGHT_GRAY)

# AI provider badges
providers = [("OpenAI", ACCENT_GREEN), ("Azure OpenAI", ACCENT_BLUE), ("Anthropic", ACCENT_PURPLE)]
for i, (name, color) in enumerate(providers):
    add_rect(slide, Inches(11 + i * 1.5), Inches(6.95), Inches(1.3), Inches(0.35), color, radius=0.1)
    add_text(slide, Inches(11 + i * 1.5), Inches(6.97), Inches(1.3), Inches(0.33),
             name, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_page_number(slide, 13, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 14 — ARQUITECTURA TECNICA
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "ARQUITECTURA", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Stack Moderna e Escalavel",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# Architecture layers
layers = [
    ("FRONTEND", Inches(1.7), [
        ("Next.js 15", ACCENT_BLUE), ("React 19", ACCENT_CYAN),
        ("Tailwind CSS", ACCENT_GREEN), ("TanStack Table", ACCENT_PURPLE),
        ("Recharts", ACCENT_AMBER), ("Radix UI", ACCENT_ROSE),
    ]),
    ("BACKEND", Inches(3.5), [
        ("API Routes", ACCENT_BLUE), ("Services Layer", ACCENT_CYAN),
        ("Repository Pattern", ACCENT_GREEN), ("TypeScript", ACCENT_PURPLE),
        ("Zod Validation", ACCENT_AMBER), ("Auth/Session", ACCENT_ROSE),
    ]),
    ("SERVICOS", Inches(5.3), [
        ("Inbox Sync", ACCENT_BLUE), ("CV Parsing", ACCENT_CYAN),
        ("Scoring Engine", ACCENT_GREEN), ("Email Service", ACCENT_PURPLE),
        ("Duplicate Detection", ACCENT_AMBER), ("Calendar Service", ACCENT_ROSE),
    ]),
    ("DADOS & STORAGE", Inches(7.1), [
        ("SQLite + Prisma", ACCENT_BLUE), ("Filesystem Local", ACCENT_CYAN),
        ("13 Modelos DB", ACCENT_GREEN), ("Audit Logging", ACCENT_PURPLE),
        ("Migrations", ACCENT_AMBER), ("Seed Data", ACCENT_ROSE),
    ]),
]

for layer_name, y, items in layers:
    add_rect(slide, Inches(0.5), y, Inches(15), Inches(1.5), CARD_BG, radius=0.03)
    add_text(slide, Inches(0.8), y + Inches(0.15), Inches(2.5), Inches(0.35),
             layer_name, font_size=13, color=ACCENT_BLUE, bold=True)

    for i, (tech, color) in enumerate(items):
        tx = Inches(0.8 + i * 2.45)
        add_rect(slide, tx, y + Inches(0.65), Inches(2.2), Inches(0.55),
                 darken(color, 5), radius=0.08)
        add_rect(slide, tx, y + Inches(0.65), Inches(0.06), Inches(0.55), color)
        add_text(slide, tx + Inches(0.2), y + Inches(0.72), Inches(1.9), Inches(0.4),
                 tech, font_size=11, color=color, bold=True)

    # Arrow down between layers (except last)
    if y < Inches(7):
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, Inches(7.7), y + Inches(1.5), Inches(0.3), Inches(0.25)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_CYAN
        arrow.line.fill.background()

add_bottom_bar(slide)
add_page_number(slide, 14, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 15 — ROADMAP
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "ROADMAP", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Evolucao Planeada em Tres Fases",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# Timeline line
add_rect(slide, Inches(1), Inches(2.7), Inches(14), Inches(0.06), ACCENT_BLUE)

phases = [
    ("FASE 1", "Fundacao ATS", "CONCLUIDO", ACCENT_GREEN,
     ["Autenticacao e sessoes", "Dashboard e metricas", "Gestao de candidatos",
      "Sync de mailbox M365", "Parsing de CVs (PDF/DOCX)", "Motor de scoring",
      "Templates de email", "Envio via Graph"]),
    ("FASE 2", "Produtividade Operacional", "EM CURSO", ACCENT_AMBER,
     ["Importacao de vagas JSON", "Ranking por vaga", "Accoes em massa",
      "Exportacao Excel", "Deteccao de duplicados", "Talent Pool inteligente"]),
    ("FASE 3", "Automatizacao Avancada", "PLANEADO", ACCENT_PURPLE,
     ["OCR para CVs digitalizados", "Integracao Calendar M365", "Resumos com IA",
      "Sugestoes de email com IA", "Relatorios avancados", "Fluxo OAuth robusto"]),
]

for i, (phase, title, status, color, items) in enumerate(phases):
    cx = Inches(0.5 + i * 5.1)

    # Timeline dot
    add_icon_circle(slide, cx + Inches(2), Inches(2.45), Inches(0.5), color, str(i+1))

    # Card
    card = add_rect(slide, cx, Inches(3.2), Inches(4.8), Inches(5.0), CARD_BG, radius=0.04)
    add_rect(slide, cx, Inches(3.2), Inches(4.8), Inches(0.06), color)

    add_text(slide, cx + Inches(0.3), Inches(3.45), Inches(1.2), Inches(0.3),
             phase, font_size=12, color=color, bold=True)

    # Status badge
    add_rect(slide, cx + Inches(2.5), Inches(3.4), Inches(2), Inches(0.35),
             darken(color), radius=0.15)
    add_text(slide, cx + Inches(2.5), Inches(3.42), Inches(2), Inches(0.33),
             status, font_size=9, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, cx + Inches(0.3), Inches(3.9), Inches(4.2), Inches(0.4),
             title, font_size=16, color=WHITE, bold=True)

    for j, item in enumerate(items):
        iy = Inches(4.5 + j * 0.42)
        add_rect(slide, cx + Inches(0.4), iy + Inches(0.05), Inches(0.15), Inches(0.15), color, radius=0.5)
        add_text(slide, cx + Inches(0.7), iy, Inches(3.8), Inches(0.3),
                 item, font_size=11, color=LIGHT_GRAY)

add_bottom_bar(slide)
add_page_number(slide, 15, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════
# SLIDE 16 — CTA / ENCERRAMENTO
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

# Decorative elements
c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-3), Inches(12), Inches(12))
c1.fill.solid()
c1.fill.fore_color.rgb = RGBColor(0x15, 0x20, 0x3C)
c1.line.fill.background()

c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(3), Inches(10), Inches(10))
c2.fill.solid()
c2.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
c2.line.fill.background()

add_gradient_bar(slide, Inches(0), Inches(0), Inches(16), Inches(0.08), ACCENT_BLUE, ACCENT_CYAN)

add_text(slide, Inches(1), Inches(2.0), Inches(14), Inches(1.0),
         "Pronto para Transformar", font_size=48, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")
add_text(slide, Inches(1), Inches(3.0), Inches(14), Inches(1.0),
         "o Recrutamento?", font_size=48, color=ACCENT_CYAN, bold=True,
         alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")

add_rect(slide, Inches(6.5), Inches(4.2), Inches(3), Inches(0.04), ACCENT_BLUE)

add_text(slide, Inches(2), Inches(4.6), Inches(12), Inches(0.8),
         "O ITSector Talent Inbox ATS esta pronto para implementacao imediata.\n"
         "100% local. Seguro. Inteligente.",
         font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key stats row
stats = [
    ("12+", "Funcionalidades\nPrincipais"),
    ("13", "Modelos de\nDados"),
    ("8", "Vagas\nPre-configuradas"),
    ("3", "Fases de\nEvolucao"),
]
for i, (val, label) in enumerate(stats):
    sx = Inches(1.5 + i * 3.5)
    add_text(slide, sx, Inches(5.8), Inches(2.5), Inches(0.6),
             val, font_size=42, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, sx, Inches(6.4), Inches(2.5), Inches(0.6),
             label, font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# CTA Button
add_rect(slide, Inches(5.5), Inches(7.3), Inches(5), Inches(0.7), ACCENT_BLUE, radius=0.1)
add_text(slide, Inches(5.5), Inches(7.35), Inches(5), Inches(0.6),
         "Comecar Agora  >>", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(3), Inches(8.2), Inches(10), Inches(0.4),
         "ITSector  |  careers@itsector.pt  |  Talent Inbox ATS",
         font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_page_number(slide, 16, TOTAL_SLIDES)


# ─── GUARDAR ─────────────────────────────────────────────────────
output_path = os.path.join(
    r"C:\Users\e617\OneDrive - ALTEN Group\Documentos\Projetos\ATs_talent",
    "ITSector_Talent_ATS_Comercial.pptx"
)
prs.save(output_path)
print(f"Apresentacao guardada em: {output_path}")
print(f"Total de slides: {TOTAL_SLIDES}")
