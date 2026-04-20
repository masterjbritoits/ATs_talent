#!/usr/bin/env python3
"""
Gera apresentação comercial premium para ITSector Talent Inbox ATS — v2
Novidades: Fase 4 completa, Rate Limiting, Pagination, Duplicate Merge UI,
           Arquitectura de Segurança, Roadmap actualizado, stats finais.
Design: Dark premium com acentos azul/ciano, tipografia moderna (mesmo que v1)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── CORES PREMIUM ───────────────────────────────────────────────
DARK_BG        = RGBColor(0x0F, 0x17, 0x2A)
DARK_BG2       = RGBColor(0x15, 0x20, 0x3C)
ACCENT_BLUE    = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_CYAN    = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_GREEN   = RGBColor(0x10, 0xB9, 0x81)
ACCENT_PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_ROSE    = RGBColor(0xF4, 0x3F, 0x5E)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xCB, 0xD5, 0xE1)
MED_GRAY       = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG        = RGBColor(0x1E, 0x29, 0x3B)
GRADIENT_START = RGBColor(0x1E, 0x40, 0xAF)
GRADIENT_END   = RGBColor(0x06, 0xB6, 0xD4)

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

# ─── HELPERS ─────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rgb_tuple(c):
    return (c[0], c[1], c[2])

def darken(c, factor=4):
    r, g, b = rgb_tuple(c)
    return RGBColor(r // factor, g // factor, b // factor)

def add_gradient_bar(slide, left, top, width, height, color1, color2):
    seg_w = width // 3
    r1,g1,b1 = rgb_tuple(color1)
    r2,g2,b2 = rgb_tuple(color2)
    colors = [color1, RGBColor((r1+r2)//2,(g1+g2)//2,(b1+b2)//2), color2]
    for i, c in enumerate(colors):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left+seg_w*i, top, seg_w+Emu(5), height)
        shape.fill.solid(); shape.fill.fore_color.rgb = c; shape.line.fill.background()

def add_rect(slide, left, top, width, height, color, radius=None):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s  = slide.shapes.add_shape(st, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    if radius: s.adjustments[0] = radius
    return s

def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font_name; p.alignment = alignment
    return tb

def add_icon_circle(slide, left, top, size, color, icon_text=""):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    if icon_text:
        tf = s.text_frame; tf.word_wrap = False
        p  = tf.paragraphs[0]; p.text = icon_text
        p.font.size = Pt(int(Emu(size)/914400*14))
        p.font.color.rgb = WHITE; p.font.bold = True
        p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
    return s

def add_metric_card(slide, left, top, width, height, value, label, color):
    add_rect(slide, left, top, width, height, CARD_BG, radius=0.05)
    add_text(slide, left, top+Inches(0.25), width, Inches(0.6),
             value, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top+Inches(0.85), width, Inches(0.4),
             label, font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

def add_page_number(slide, num, total):
    add_text(slide, Inches(15.0), Inches(8.55), Inches(0.8), Inches(0.3),
             f"{num}/{total}", font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)

def add_bottom_bar(slide):
    add_gradient_bar(slide, Inches(0), Inches(8.85), Inches(16), Inches(0.15), ACCENT_BLUE, ACCENT_CYAN)

def add_sidebar(slide, active_index):
    items = ["Dashboard","Inbox","Candidatos","Vagas","Talent Pool","Templates","Relatorios","Definicoes"]
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.5), Inches(7.0), RGBColor(0x0F,0x17,0x2A))
    add_text(slide, Inches(0.6), Inches(1.6), Inches(2.2), Inches(0.4),
             "ITSector ATS", font_size=13, color=ACCENT_CYAN, bold=True)
    for i, name in enumerate(items):
        y = Inches(2.0 + i*0.45)
        if i == active_index:
            add_rect(slide, Inches(0.5), y, Inches(2.5), Inches(0.4), ACCENT_BLUE, radius=0.08)
        add_text(slide, Inches(0.9), y+Inches(0.02), Inches(2), Inches(0.35),
                 name, font_size=11, color=WHITE if i==active_index else MED_GRAY)

TOTAL_SLIDES = 17

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA  (v2: badge "v2.0 | Fase 4 Completa", novas métricas)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(10), Inches(10))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x1E,0x29,0x3B); c1.line.fill.background()
c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(4), Inches(8), Inches(8))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x15,0x20,0x3C); c2.line.fill.background()

add_gradient_bar(slide, Inches(0), Inches(0), Inches(16), Inches(0.08), ACCENT_BLUE, ACCENT_CYAN)

# Badge v2.0
add_rect(slide, Inches(1.5), Inches(2.2), Inches(3.6), Inches(0.45), ACCENT_BLUE, radius=0.15)
add_text(slide, Inches(1.5), Inches(2.22), Inches(3.6), Inches(0.45),
         "PLATAFORMA ATS  |  v2.0  |  Fase 4 Completa",
         font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
         font_name="Segoe UI Semibold")

add_text(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(1.2),
         "ITSector Talent", font_size=54, color=WHITE, bold=True, font_name="Segoe UI Light")
add_text(slide, Inches(1.5), Inches(3.9), Inches(10), Inches(1.2),
         "Inbox ATS", font_size=54, color=ACCENT_CYAN, bold=True, font_name="Segoe UI Light")

add_text(slide, Inches(1.5), Inches(5.2), Inches(8), Inches(0.8),
         "Sistema Inteligente de Gestao de Candidaturas e Recrutamento",
         font_size=20, color=LIGHT_GRAY, font_name="Segoe UI")
add_rect(slide, Inches(1.5), Inches(6.1), Inches(3), Inches(0.04), ACCENT_BLUE)
add_text(slide, Inches(1.5), Inches(6.4), Inches(6), Inches(0.5),
         "Transformar emails em talento. Automaticamente.",
         font_size=16, color=MED_GRAY, font_name="Segoe UI")

# Mini-métricas actualizadas para v2
for i, (val, lbl) in enumerate([
    ("19+", "Funcionalidades"),
    ("M365", "Integracao Nativa"),
    ("100%", "Local & Seguro"),
]):
    x = Inches(1.5 + i * 3.2)
    add_text(slide, x, Inches(7.3),  Inches(2), Inches(0.4), val, font_size=22, color=ACCENT_CYAN, bold=True)
    add_text(slide, x, Inches(7.65), Inches(2.5), Inches(0.3), lbl, font_size=12, color=MED_GRAY)

add_bottom_bar(slide); add_page_number(slide, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — O DESAFIO  (igual ao v1)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.5), Inches(8), Inches(0.6),
         "O DESAFIO", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(1.0), Inches(12), Inches(0.8),
         "O Recrutamento Tradicional por Email Nao Escala",
         font_size=36, color=WHITE, bold=True, font_name="Segoe UI Light")

pains = [
    ("Emails Perdidos",    "Candidaturas enterradas na caixa\nde entrada sem rastreamento",    ACCENT_ROSE,   "!"),
    ("CVs Manuais",        "Leitura manual de dezenas de CVs\npor dia sem normalizacao",        ACCENT_AMBER,  "?"),
    ("Sem Ranking",        "Impossivel comparar candidatos\nobjectivamente entre vagas",        ACCENT_PURPLE, "#"),
    ("Tempo Perdido",      "Horas gastas em tarefas repetitivas\nque podiam ser automatizadas", ACCENT_ROSE,   "T"),
]
for i, (title, desc, color, icon) in enumerate(pains):
    x = Inches(1 + i*3.6); y = Inches(2.5)
    add_rect(slide, x, y, Inches(3.2), Inches(3.5), CARD_BG, radius=0.05)
    add_rect(slide, x, y, Inches(3.2), Inches(0.06), color)
    add_icon_circle(slide, x+Inches(1.2), y+Inches(0.4), Inches(0.7), color, icon)
    add_text(slide, x+Inches(0.3), y+Inches(1.3), Inches(2.6), Inches(0.5),
             title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x+Inches(0.3), y+Inches(1.9), Inches(2.6), Inches(1.2),
             desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(1), Inches(6.8), Inches(14), Inches(1.2), DARK_BG2, radius=0.03)
add_text(slide, Inches(1.5), Inches(7.0),  Inches(2),  Inches(0.8), "67%", font_size=42, color=ACCENT_CYAN, bold=True)
add_text(slide, Inches(3.5), Inches(7.05), Inches(11), Inches(0.8),
         "dos recrutadores afirmam perder candidatos qualificados devido a processos manuais desorganizados",
         font_size=16, color=LIGHT_GRAY)

add_bottom_bar(slide); add_page_number(slide, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — A SOLUCAO
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.5), Inches(8), Inches(0.6),
         "A SOLUCAO", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(1.0), Inches(13), Inches(0.8),
         "De Email a Contratacao - Tudo Automatizado",
         font_size=36, color=WHITE, bold=True, font_name="Segoe UI Light")

steps = [
    ("1", "Email Recebido",  "Sincronizacao\nautomatica M365",  ACCENT_BLUE),
    ("2", "CV Processado",   "Parsing PDF/DOCX\ne extracao de dados", ACCENT_CYAN),
    ("3", "Candidato Criado","Perfil normalizado\ncom competencias",   ACCENT_GREEN),
    ("4", "Score Calculado", "Avaliacao heuristica\npor vaga",         ACCENT_PURPLE),
    ("5", "Decisao Tomada",  "Avancar, Rever\nou Rejeitar",            ACCENT_AMBER),
]
for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i*3.05); y = Inches(2.5)
    add_rect(slide, x, y, Inches(2.7), Inches(3.0), CARD_BG, radius=0.05)
    add_rect(slide, x, y, Inches(2.7), Inches(0.06), color)
    add_icon_circle(slide, x+Inches(0.95), y+Inches(0.3), Inches(0.7), color, num)
    add_text(slide, x+Inches(0.15), y+Inches(1.2), Inches(2.4), Inches(0.45),
             title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x+Inches(0.15), y+Inches(1.7), Inches(2.4), Inches(0.8),
             desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x+Inches(2.7), y+Inches(1.2), Inches(0.35), Inches(0.3))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT_CYAN; arrow.line.fill.background()

add_rect(slide, Inches(1), Inches(6.3), Inches(14), Inches(1.8), DARK_BG2, radius=0.03)
add_text(slide, Inches(1.5), Inches(6.5), Inches(13), Inches(0.5),
         "Resultado: Zero Trabalho Manual na Triagem Inicial",
         font_size=22, color=ACCENT_CYAN, bold=True)
add_text(slide, Inches(1.5), Inches(7.1), Inches(13), Inches(0.8),
         "O sistema processa automaticamente cada email, extrai informacao dos CVs, cria perfis de candidatos,\n"
         "calcula pontuacoes de adequacao e recomenda accoes - tudo em segundos.",
         font_size=14, color=LIGHT_GRAY)

add_bottom_bar(slide); add_page_number(slide, 3, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — DASHBOARD  (v2: pipeline mais populado, badge v2)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "DASHBOARD", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(10), Inches(0.6),
         "Visao Geral em Tempo Real", font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

add_rect(slide, Inches(0.5), Inches(1.5), Inches(15), Inches(7.0), RGBColor(0x0B,0x11,0x1F), radius=0.02)
add_sidebar(slide, 0)

for i, (val, label, color) in enumerate([
    ("312", "Total Candidatos", ACCENT_BLUE),
    ("14",  "Vagas Abertas",    ACCENT_GREEN),
    ("41",  "Em Revisao",       ACCENT_AMBER),
    ("92%", "Taxa Processamento", ACCENT_CYAN),
]):
    x = Inches(3.3 + i*2.85)
    add_metric_card(slide, x, Inches(1.8), Inches(2.6), Inches(1.5), val, label, color)

# Gráfico barras
add_rect(slide, Inches(3.3), Inches(3.6), Inches(5.5), Inches(3.2), CARD_BG, radius=0.03)
add_text(slide, Inches(3.6), Inches(3.7), Inches(4), Inches(0.4),
         "Candidaturas por Vaga", font_size=12, color=LIGHT_GRAY, bold=True)
for i, (label, pct, color) in enumerate([
    (".NET Dev",    0.90, ACCENT_BLUE),
    ("React Dev",   0.72, ACCENT_CYAN),
    ("QA Engineer", 0.48, ACCENT_GREEN),
    ("DevOps",      0.58, ACCENT_PURPLE),
    ("Power Apps",  0.38, ACCENT_AMBER),
    ("Mobile Dev",  0.44, ACCENT_ROSE),
]):
    by = Inches(4.25 + i*0.42)
    add_text(slide, Inches(3.6), by, Inches(1.2), Inches(0.3), label, font_size=9, color=MED_GRAY)
    add_rect(slide, Inches(5.0), by+Inches(0.04), Inches(3.5*pct), Inches(0.22), color, radius=0.2)

# Actividade recente
add_rect(slide, Inches(9.1), Inches(3.6), Inches(6.1), Inches(3.2), CARD_BG, radius=0.03)
add_text(slide, Inches(9.4), Inches(3.7), Inches(4), Inches(0.4),
         "Actividade Recente", font_size=12, color=LIGHT_GRAY, bold=True)
for i, act in enumerate([
    "Duplicado fundido: Joao Santos (2 registos)",
    "Candidato criado: Ana Rodrigues (React Dev)",
    "Score: Pedro Almeida - 91/100 -> AVANCAR",
    "Entrevista agendada: Maria Silva",
    "Exportacao Excel: 312 candidatos",
    "Vaga publicada: DevOps Engineer",
]):
    add_rect(slide, Inches(9.3), Inches(4.2+i*0.42), Inches(0.12), Inches(0.12),
             [ACCENT_AMBER,ACCENT_BLUE,ACCENT_GREEN,ACCENT_PURPLE,ACCENT_CYAN,ACCENT_ROSE][i])
    add_text(slide, Inches(9.6), Inches(4.15+i*0.42), Inches(5.3), Inches(0.3),
             act, font_size=10, color=LIGHT_GRAY)

# Pipeline mini
add_rect(slide, Inches(3.3), Inches(7.0), Inches(11.9), Inches(1.2), CARD_BG, radius=0.03)
add_text(slide, Inches(3.6), Inches(7.05), Inches(4), Inches(0.35),
         "Pipeline de Recrutamento", font_size=12, color=LIGHT_GRAY, bold=True)
for i, (name, count, color) in enumerate([
    ("Novo",      "52", ACCENT_BLUE),
    ("Revisao",   "31", ACCENT_AMBER),
    ("Shortlist", "18", ACCENT_CYAN),
    ("Entrevista","10", ACCENT_PURPLE),
    ("Oferta",    "4",  ACCENT_GREEN),
    ("Contratado","6",  ACCENT_GREEN),
]):
    sx = Inches(3.5 + i*1.95)
    add_rect(slide, sx, Inches(7.45), Inches(1.7), Inches(0.6), color, radius=0.05)
    add_text(slide, sx, Inches(7.45), Inches(1.7), Inches(0.3),
             count, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, sx, Inches(7.73), Inches(1.7), Inches(0.25),
             name, font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide); add_page_number(slide, 4, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — INBOX & SYNC
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "INBOX INTELIGENTE", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Sincronizacao Automatica com Microsoft 365",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

add_rect(slide, Inches(0.5), Inches(1.5), Inches(15), Inches(6.8), RGBColor(0x0B,0x11,0x1F), radius=0.02)
add_sidebar(slide, 1)

add_text(slide, Inches(3.3), Inches(1.7), Inches(5), Inches(0.4),
         "Inbox - careers@itsector.pt", font_size=16, color=WHITE, bold=True)
add_rect(slide, Inches(12.5), Inches(1.65), Inches(2.5), Inches(0.5), ACCENT_BLUE, radius=0.1)
add_text(slide, Inches(12.5), Inches(1.68), Inches(2.5), Inches(0.45),
         "Sincronizar Agora", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(3.3), Inches(2.3), Inches(12), Inches(0.45), DARK_BG2, radius=0.05)
add_text(slide, Inches(3.6), Inches(2.32), Inches(10), Inches(0.35),
         "Ultima sincronizacao: ha 3 minutos  |  29 emails processados  |  11 novos candidatos criados",
         font_size=10, color=ACCENT_GREEN)

for i, (name, subject, time, status, status_color, is_cand) in enumerate([
    ("Maria Silva",    "Candidatura - .NET Developer Senior",  "14:32", "CV Processado", ACCENT_GREEN, True),
    ("Joao Santos",    "Candidatura React Developer",          "13:45", "Em Analise",    ACCENT_AMBER, True),
    ("Ana Costa",      "Re: Oportunidade QA Engineer",         "12:20", "Novo",          ACCENT_BLUE,  True),
    ("Pedro Almeida",  "Candidatura Espontanea - DevOps",      "11:05", "CV Processado", ACCENT_GREEN, True),
    ("Newsletter SAPO","Tech News Semanal",                    "10:30", "Ignorado",      MED_GRAY,     False),
    ("Sofia Martins",  "Interesse vaga Power Apps",            "09:15", "Novo",          ACCENT_BLUE,  True),
    ("Carlos Ferreira","CV em anexo - Mobile Developer",       "08:40", "CV Processado", ACCENT_GREEN, True),
]):
    ey = Inches(2.95 + i*0.65)
    add_rect(slide, Inches(3.3), ey, Inches(12), Inches(0.6),
             CARD_BG if i%2==0 else RGBColor(0x17,0x22,0x34), radius=0.02)
    if is_cand:
        add_rect(slide, Inches(3.3), ey, Inches(0.06), Inches(0.6), ACCENT_GREEN)
    add_text(slide, Inches(3.6),  ey+Inches(0.05), Inches(2),   Inches(0.25), name,    font_size=11, color=WHITE, bold=True)
    add_text(slide, Inches(5.6),  ey+Inches(0.05), Inches(5.5), Inches(0.25), subject, font_size=11, color=LIGHT_GRAY)
    add_text(slide, Inches(12.5), ey+Inches(0.05), Inches(0.8), Inches(0.25), time,    font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(13.5), ey+Inches(0.1), Inches(1.3), Inches(0.32), darken(status_color), radius=0.15)
    add_text(slide, Inches(13.5), ey+Inches(0.12), Inches(1.3), Inches(0.3),
             status, font_size=9, color=status_color, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(3.5), Inches(7.75), Inches(3), Inches(0.3),
         "Verde = Candidatura Detectada", font_size=10, color=ACCENT_GREEN)
add_text(slide, Inches(7),   Inches(7.75), Inches(3), Inches(0.3),
         "Cinza = Email Ignorado (nao-candidatura)", font_size=10, color=MED_GRAY)

add_bottom_bar(slide); add_page_number(slide, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — CANDIDATOS  (v2: paginacao + alert de duplicado visivel)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "GESTAO DE CANDIDATOS", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Pipeline com Filtros, Paginacao e Revisao de Duplicados",
         font_size=28, color=WHITE, bold=True, font_name="Segoe UI Light")

add_rect(slide, Inches(0.5), Inches(1.5), Inches(15), Inches(6.8), RGBColor(0x0B,0x11,0x1F), radius=0.02)
add_sidebar(slide, 2)

# Topo da tabela com filtros
add_text(slide, Inches(3.3), Inches(1.65), Inches(5), Inches(0.4),
         "Candidate Pipeline", font_size=15, color=WHITE, bold=True)
add_rect(slide, Inches(12.5), Inches(1.6), Inches(2.5), Inches(0.45), ACCENT_BLUE, radius=0.1)
add_text(slide, Inches(12.5), Inches(1.63), Inches(2.5), Inches(0.43),
         "Accoes em Massa", font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Filtros
add_rect(slide, Inches(3.3), Inches(2.15), Inches(12), Inches(0.4), DARK_BG2, radius=0.04)
for i, f in enumerate(["Status: Todos","Vaga: Todas","Score: 0+","Ordenar: Score"]):
    add_rect(slide, Inches(3.5+i*2.9), Inches(2.22), Inches(2.5), Inches(0.26), CARD_BG, radius=0.08)
    add_text(slide, Inches(3.6+i*2.9), Inches(2.23), Inches(2.3), Inches(0.25),
             f, font_size=9, color=LIGHT_GRAY)

# Header tabela
add_rect(slide, Inches(3.3), Inches(2.65), Inches(12), Inches(0.35), ACCENT_BLUE, radius=0.03)
for i, h in enumerate(["#","Nome","Vaga","Score","Status","Accao"]):
    widths = [0.3, 2.5, 2.5, 0.9, 1.2, 1.2]
    x = Inches(3.4 + sum(widths[:i]))
    add_text(slide, x, Inches(2.67), Inches(widths[i]), Inches(0.3),
             h, font_size=9, color=WHITE, bold=True)

# Linhas candidatos
rows = [
    (1, "Maria Silva",    ".NET Developer",   "91", "SHORTLISTED",  ACCENT_CYAN,  ACCENT_GREEN),
    (2, "Pedro Almeida",  "DevOps Engineer",  "84", "SHORTLISTED",  ACCENT_CYAN,  ACCENT_GREEN),
    (3, "Joao Santos",    "React Developer",  "76", "MANUAL_REVIEW",ACCENT_AMBER, ACCENT_AMBER),
    (4, "Ana Costa",      "QA Engineer",      "68", "MANUAL_REVIEW",ACCENT_AMBER, ACCENT_AMBER),
    (5, "Sofia Martins",  "Power Apps Dev",   "59", "NEW",          ACCENT_BLUE,  ACCENT_BLUE),
    (6, "Carlos Ferreira","Mobile Developer", "52", "NEW",          ACCENT_BLUE,  ACCENT_BLUE),
]
for i, (num, name, job, score, status, s_color, btn_color) in enumerate(rows):
    ry = Inches(3.05 + i*0.48)
    add_rect(slide, Inches(3.3), ry, Inches(12), Inches(0.44),
             CARD_BG if i%2==0 else RGBColor(0x17,0x22,0x34), radius=0.02)
    # checkbox
    add_rect(slide, Inches(3.4), ry+Inches(0.12), Inches(0.2), Inches(0.2), DARK_BG2, radius=0.1)
    add_text(slide, Inches(3.65), ry+Inches(0.06), Inches(0.3), Inches(0.3), str(num),   font_size=10, color=MED_GRAY)
    add_text(slide, Inches(3.95), ry+Inches(0.06), Inches(2.5), Inches(0.3), name,       font_size=11, color=WHITE, bold=True)
    add_text(slide, Inches(6.45), ry+Inches(0.06), Inches(2.5), Inches(0.3), job,        font_size=10, color=LIGHT_GRAY)
    add_rect(slide, Inches(8.9), ry+Inches(0.07), Inches(0.75), Inches(0.28), darken(btn_color,5), radius=0.1)
    add_text(slide, Inches(8.9), ry+Inches(0.09), Inches(0.75), Inches(0.26), score, font_size=10, color=btn_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide, Inches(9.7), ry+Inches(0.07), Inches(1.1), Inches(0.28), darken(s_color,5), radius=0.1)
    add_text(slide, Inches(9.7), ry+Inches(0.09), Inches(1.1), Inches(0.26),
             status, font_size=8, color=s_color, bold=True, alignment=PP_ALIGN.CENTER)

# ── Paginação (nova feature v2) ──
add_rect(slide, Inches(3.3), Inches(5.95), Inches(12), Inches(0.45), DARK_BG2, radius=0.03)
add_text(slide, Inches(3.6), Inches(6.0), Inches(4), Inches(0.35),
         "312 candidatos no total", font_size=10, color=MED_GRAY)
for i, lbl in enumerate(["<< 1","< 1","1","2 >","3 >>>"]):
    active = lbl == "1"
    bx = Inches(10.5 + i*0.55)
    add_rect(slide, bx, Inches(6.0), Inches(0.48), Inches(0.32),
             ACCENT_BLUE if active else DARK_BG, radius=0.08)
    add_text(slide, bx, Inches(6.01), Inches(0.48), Inches(0.3),
             lbl, font_size=9, color=WHITE if active else MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── Alert duplicados (nova feature v2) ──
add_rect(slide, Inches(3.3), Inches(6.55), Inches(12), Inches(0.65),
         RGBColor(0x2D, 0x1F, 0x00), radius=0.04)
add_rect(slide, Inches(3.3), Inches(6.55), Inches(0.06), Inches(0.65), ACCENT_AMBER)
add_text(slide, Inches(3.5), Inches(6.6), Inches(7), Inches(0.35),
         "! Fila de Revisao de Duplicados  —  3 grupos detectados",
         font_size=12, color=ACCENT_AMBER, bold=True)
add_text(slide, Inches(3.5), Inches(6.9), Inches(7), Inches(0.25),
         "Confianca 90%+  |  Telefone e email coincidentes  |  Expandir para fundir",
         font_size=10, color=MED_GRAY)
add_rect(slide, Inches(14.2), Inches(6.65), Inches(0.85), Inches(0.4), ACCENT_AMBER, radius=0.1)
add_text(slide, Inches(14.2), Inches(6.67), Inches(0.85), Inches(0.38),
         "Rever", font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide); add_page_number(slide, 6, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — CV PARSING
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "PROCESSAMENTO DE CVs", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Extracao Inteligente de Curriculos",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# CV doc
add_rect(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(6.5), RGBColor(0xF8,0xFA,0xFC), radius=0.03)
add_text(slide, Inches(1.2), Inches(1.9), Inches(4), Inches(0.4),
         "curriculum_vitae.pdf", font_size=14, color=RGBColor(0x1E,0x29,0x3B), bold=True)
add_rect(slide, Inches(1.2), Inches(2.35), Inches(4.2), Inches(0.04), RGBColor(0xE2,0xE8,0xF0))
for i, line in enumerate([
    "MARIA SILVA","Senior .NET Developer","",
    "EXPERIENCIA PROFISSIONAL",
    "- ITSector (2019-presente) - .NET Developer",
    "- BNP Paribas (2016-2019) - Software Engineer","",
    "COMPETENCIAS TECNICAS",
    "C#, .NET Core, Azure, SQL Server, Docker",
    "REST APIs, Microservices, CI/CD, Git","",
    "IDIOMAS","Portugues (Nativo), Ingles (Fluente)","",
    "FORMACAO","MSc Engenharia Informatica - IST (2016)",
]):
    is_hdr = line.isupper() and line != ""
    add_text(slide, Inches(1.3), Inches(2.5+i*0.3), Inches(4), Inches(0.28),
             line, font_size=9 if not is_hdr else 10,
             color=RGBColor(0x0F,0x17,0x2A) if is_hdr else RGBColor(0x47,0x55,0x69),
             bold=is_hdr)

arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.0), Inches(4.2), Inches(1.2), Inches(0.6))
arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT_CYAN; arrow.line.fill.background()
add_text(slide, Inches(6.0), Inches(4.9), Inches(1.2), Inches(0.3),
         "Parsing", font_size=10, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Dados extraídos
add_rect(slide, Inches(7.5), Inches(1.7), Inches(7.5), Inches(6.5), CARD_BG, radius=0.03)
add_text(slide, Inches(7.8), Inches(1.9), Inches(5), Inches(0.4),
         "Dados Extraidos Automaticamente", font_size=14, color=ACCENT_GREEN, bold=True)
for i, (key, val, color) in enumerate([
    ("Nome",              "Maria Silva",                                    ACCENT_CYAN),
    ("Titulo Actual",     "Senior .NET Developer",                          ACCENT_CYAN),
    ("Localizacao",       "Lisboa, Portugal",                               ACCENT_CYAN),
    ("Anos Experiencia",  "8 anos",                                         ACCENT_CYAN),
    ("Skills",            '["C#", ".NET Core", "Azure", "Docker"]',         ACCENT_GREEN),
    ("Idiomas",           '["Portugues (Nativo)", "Ingles (Fluente)"]',      ACCENT_GREEN),
    ("Educacao",          '["MSc Eng. Informatica - IST 2016"]',             ACCENT_GREEN),
    ("Experiencia",       '["ITSector 2019-2025", "BNP Paribas 2016-19"]', ACCENT_GREEN),
    ("LinkedIn",          "linkedin.com/in/maria-silva-dev",                ACCENT_BLUE),
    ("Email",             "maria.silva@gmail.com",                          ACCENT_BLUE),
    ("Confianca Parser",  "94%",                                            ACCENT_GREEN),
]):
    ey = Inches(2.5 + i*0.47)
    add_text(slide, Inches(7.9), ey, Inches(2.2), Inches(0.3), key, font_size=10, color=MED_GRAY, bold=True)
    add_text(slide, Inches(10.2), ey, Inches(4.5), Inches(0.3), val, font_size=10, color=color)

add_rect(slide, Inches(7.8), Inches(7.2), Inches(6.8), Inches(0.7), DARK_BG2, radius=0.05)
for i, (fmt, clr) in enumerate([("PDF",ACCENT_BLUE),("DOCX",ACCENT_GREEN),("DOC",ACCENT_CYAN),("OCR",ACCENT_PURPLE)]):
    fx = Inches(8.2 + i*1.6)
    add_rect(slide, fx, Inches(7.35), Inches(1.2), Inches(0.4), clr, radius=0.1)
    add_text(slide, fx, Inches(7.37), Inches(1.2), Inches(0.38),
             fmt, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide); add_page_number(slide, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — SCORING ENGINE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "MOTOR DE SCORING", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Avaliacao Objectiva e Configuravel",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

cx, cy = Inches(3.5), Inches(4.5)
add_icon_circle(slide, cx-Inches(1.5), cy-Inches(1.5), Inches(3), ACCENT_BLUE, "")
add_icon_circle(slide, cx-Inches(1.2), cy-Inches(1.2), Inches(2.4), DARK_BG, "")
add_text(slide, cx-Inches(0.8), cy-Inches(0.5), Inches(1.6), Inches(0.6),
         "100", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, cx-Inches(0.5), cy+Inches(0.3), Inches(1), Inches(0.3),
         "pontos", font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

for i, (name, pts, desc, color, pct) in enumerate([
    ("Skills Obrigatorias", "30 pts", "Match com requisitos da vaga",   ACCENT_BLUE,   0.30),
    ("Anos Experiencia",    "15 pts", "Adequacao ao nivel exigido",      ACCENT_CYAN,   0.15),
    ("Skills Opcionais",    "10 pts", "Competencias diferenciadoras",    ACCENT_GREEN,  0.10),
    ("Idiomas",             "10 pts", "Correspondencia linguistica",     ACCENT_PURPLE, 0.10),
    ("Titulo/Senioridade",  "10 pts", "Relevancia funcional",            ACCENT_AMBER,  0.10),
    ("Dominio",             "10 pts", "Experiencia no sector",           ACCENT_ROSE,   0.10),
    ("Localizacao",         "5 pts",  "Proximidade geografica",          ACCENT_CYAN,   0.05),
    ("Qualidade CV",        "5 pts",  "Completude do curriculo",         MED_GRAY,      0.05),
    ("Factor Recrutador",   "5 pts",  "Ajuste manual do recruiter",      ACCENT_BLUE,   0.05),
]):
    cy2 = Inches(1.8 + i*0.65)
    add_rect(slide, Inches(7), cy2, Inches(0.4), Inches(0.4), color, radius=0.15)
    add_text(slide, Inches(7.55), cy2, Inches(3), Inches(0.25), name, font_size=12, color=WHITE, bold=True)
    add_text(slide, Inches(7.55), cy2+Inches(0.22), Inches(4), Inches(0.2), desc, font_size=9, color=MED_GRAY)
    add_rect(slide, Inches(11.5), cy2+Inches(0.08), Inches(3), Inches(0.14), DARK_BG2, radius=0.3)
    add_rect(slide, Inches(11.5), cy2+Inches(0.08), Inches(3*pct/0.30), Inches(0.14), color, radius=0.3)
    add_text(slide, Inches(14.6), cy2, Inches(0.8), Inches(0.3), pts, font_size=11, color=color, bold=True, alignment=PP_ALIGN.RIGHT)

add_rect(slide, Inches(0.8), Inches(7.3), Inches(14.4), Inches(0.9), DARK_BG2, radius=0.03)
for i, (score, action, desc, color) in enumerate([
    ("75+", "AVANCAR",        "Auto-avanco para entrevista",         ACCENT_GREEN),
    ("45-74","REVISAO MANUAL","Requer decisao do recrutador",        ACCENT_AMBER),
    ("<45",  "REJEITAR",      "Abaixo do minimo para a vaga",       ACCENT_ROSE),
]):
    tx = Inches(1.0 + i*4.8)
    add_rect(slide, tx, Inches(7.4), Inches(0.8), Inches(0.6), color, radius=0.1)
    add_text(slide, tx, Inches(7.43), Inches(0.8), Inches(0.55),
             score, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, tx+Inches(1.0), Inches(7.4),  Inches(2.5), Inches(0.3), action, font_size=12, color=color, bold=True)
    add_text(slide, tx+Inches(1.0), Inches(7.68), Inches(3.5), Inches(0.25), desc,   font_size=10, color=MED_GRAY)

add_bottom_bar(slide); add_page_number(slide, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — GESTAO DE VAGAS
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "GESTAO DE VAGAS", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Vagas com Ranking Automatico de Candidatos",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

for i, (title, dept, loc, status, applicants, sc) in enumerate([
    (".NET Developer",        "Tecnologia",   "Lisboa", "ABERTA",   22, ACCENT_GREEN),
    ("Senior .NET Developer", "Tecnologia",   "Porto",  "ABERTA",   15, ACCENT_GREEN),
    ("React Developer",       "Frontend",     "Remoto", "ABERTA",   18, ACCENT_GREEN),
    ("Power Apps Developer",  "Low Code",     "Lisboa", "ABERTA",    8, ACCENT_GREEN),
    ("Business Analyst",      "Consultoria",  "Lisboa", "ABERTA",   11, ACCENT_GREEN),
    ("QA Engineer",           "Qualidade",    "Porto",  "PAUSADA",   9, ACCENT_AMBER),
    ("DevOps Engineer",       "Infra",        "Remoto", "ABERTA",   14, ACCENT_GREEN),
    ("Mobile Developer",      "Mobile",       "Lisboa", "RASCUNHO",  4, MED_GRAY),
]):
    row = i // 4; col = i % 4
    cx = Inches(0.5+col*3.85); cy = Inches(1.7+row*3.3)
    add_rect(slide, cx, cy, Inches(3.6), Inches(3.0), CARD_BG, radius=0.04)
    add_rect(slide, cx, cy, Inches(3.6), Inches(0.05), sc)
    add_text(slide, cx+Inches(0.3), cy+Inches(0.2),  Inches(3), Inches(0.35), title,         font_size=14, color=WHITE, bold=True)
    add_text(slide, cx+Inches(0.3), cy+Inches(0.55), Inches(3), Inches(0.25), f"{dept} | {loc}", font_size=10, color=MED_GRAY)
    add_rect(slide, cx+Inches(0.3), cy+Inches(0.9), Inches(1.0), Inches(0.28), darken(sc), radius=0.15)
    add_text(slide, cx+Inches(0.3), cy+Inches(0.91), Inches(1.0), Inches(0.27),
             status, font_size=8, color=sc, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, cx+Inches(0.3), cy+Inches(1.4), Inches(1), Inches(0.5),
             str(applicants), font_size=28, color=ACCENT_CYAN, bold=True)
    add_text(slide, cx+Inches(1.3), cy+Inches(1.5), Inches(2), Inches(0.3), "candidatos", font_size=11, color=MED_GRAY)
    for j in range(min(3, applicants)):
        bp = [0.95, 0.78, 0.62][j]; by = cy+Inches(2.0+j*0.28)
        add_text(slide, cx+Inches(0.3), by, Inches(0.3), Inches(0.2), f"#{j+1}", font_size=8, color=MED_GRAY)
        add_rect(slide, cx+Inches(0.65), by+Inches(0.03), Inches(2.2), Inches(0.14), DARK_BG2, radius=0.3)
        add_rect(slide, cx+Inches(0.65), by+Inches(0.03), Inches(2.2*bp), Inches(0.14),
                 [ACCENT_GREEN,ACCENT_CYAN,ACCENT_AMBER][j], radius=0.3)

add_bottom_bar(slide); add_page_number(slide, 9, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — EMAIL & COMUNICACAO
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "COMUNICACAO INTEGRADA", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Templates, Emails e Agendamento de Entrevistas",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

# Templates
add_rect(slide, Inches(0.5), Inches(1.6), Inches(7.2), Inches(6.5), CARD_BG, radius=0.03)
add_text(slide, Inches(0.9), Inches(1.8), Inches(5), Inches(0.4),
         "Templates de Email", font_size=16, color=WHITE, bold=True)
for i, (name, desc, color) in enumerate([
    ("Convite para Entrevista",   "Agenda automaticamente via M365 Calendar",    ACCENT_PURPLE),
    ("Rejeicao Profissional",     "Resposta respeitosa com feedback",             ACCENT_ROSE),
    ("Proposta/Oferta",           "Template com detalhes da posicao",             ACCENT_GREEN),
    ("Pedido de Informacao",      "Solicitar dados adicionais ao candidato",      ACCENT_AMBER),
    ("Boas-vindas Talent Pool",   "Informar inclusao no talent pool",             ACCENT_CYAN),
]):
    ty = Inches(2.5 + i*0.95)
    add_rect(slide, Inches(0.8), ty, Inches(6.6), Inches(0.8), DARK_BG2, radius=0.03)
    add_rect(slide, Inches(0.8), ty, Inches(0.08), Inches(0.8), color)
    add_text(slide, Inches(1.2), ty+Inches(0.08), Inches(5), Inches(0.3), name, font_size=12, color=WHITE, bold=True)
    add_text(slide, Inches(1.2), ty+Inches(0.38), Inches(5.5), Inches(0.3), desc, font_size=10, color=MED_GRAY)
    add_rect(slide, Inches(6.2), ty+Inches(0.2), Inches(1), Inches(0.35), color, radius=0.1)
    add_text(slide, Inches(6.2), ty+Inches(0.22), Inches(1), Inches(0.33),
             "Usar", font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Composer
add_rect(slide, Inches(8.2), Inches(1.6), Inches(7.2), Inches(6.5), CARD_BG, radius=0.03)
add_text(slide, Inches(8.5), Inches(1.8), Inches(5), Inches(0.4), "Compor Email", font_size=16, color=WHITE, bold=True)
for i, (label, value) in enumerate([
    ("Para:",     "maria.silva@gmail.com"),
    ("Assunto:",  "Convite para Entrevista - .NET Developer"),
]):
    fy = Inches(2.5 + i*0.6)
    add_text(slide, Inches(8.5), fy, Inches(1), Inches(0.3), label, font_size=11, color=MED_GRAY)
    add_rect(slide, Inches(9.5), fy, Inches(5.5), Inches(0.4), DARK_BG2, radius=0.05)
    add_text(slide, Inches(9.7), fy+Inches(0.03), Inches(5), Inches(0.3), value, font_size=11, color=LIGHT_GRAY)

add_rect(slide, Inches(8.5), Inches(3.9), Inches(6.6), Inches(3.0), DARK_BG2, radius=0.05)
for i, line in enumerate([
    "Cara Maria Silva,","",
    "Temos o prazer de a convidar para uma","entrevista relativa a posicao de .NET Developer","na ITSector.","",
    "Data: 15 de Maio de 2026, 14:30","Local: Escritorio ITSector Lisboa","",
    "Aguardamos a sua confirmacao.","Com os melhores cumprimentos,","Equipa de Recrutamento ITSector",
]):
    add_text(slide, Inches(8.7), Inches(4.0+i*0.24), Inches(6), Inches(0.22),
             line, font_size=10, color=LIGHT_GRAY)

add_rect(slide, Inches(13.2), Inches(7.2), Inches(1.8), Inches(0.5), ACCENT_BLUE,   radius=0.1)
add_text(slide, Inches(13.2), Inches(7.23), Inches(1.8), Inches(0.45),
         "Enviar via M365", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(11.2), Inches(7.2), Inches(1.8), Inches(0.5), ACCENT_PURPLE, radius=0.1)
add_text(slide, Inches(11.2), Inches(7.23), Inches(1.8), Inches(0.45),
         "Agendar M365", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide); add_page_number(slide, 10, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — TALENT POOL
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "TALENT POOL", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Base de Talento Reutilizavel para Futuras Oportunidades",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")
add_text(slide, Inches(1), Inches(1.5), Inches(13), Inches(0.5),
         "Candidatos qualificados que nao encaixam na vaga actual sao preservados para futuras oportunidades.",
         font_size=14, color=LIGHT_GRAY)

for i, (name, role, score_info, skills, suggestion, color) in enumerate([
    ("Ana Rodrigues", "React Developer Junior", "Score: 58 (Near-Fit QA)",
     ["React","TypeScript","Node.js"], "Sugestao: React Developer (quando abrir junior)", ACCENT_CYAN),
    ("Bruno Mendes",  "DevOps Engineer",        "Score: 62 (Near-Fit)",
     ["AWS","Terraform","K8s"],      "Sugestao: Proxima vaga Cloud/Infra",              ACCENT_PURPLE),
    ("Carla Dias",    "Business Analyst",        "Score: 71 (Revisao Manual)",
     ["BPMN","SQL","Power BI"],      "Sugestao: Rever quando abrir BA Senior",           ACCENT_AMBER),
    ("Daniel Sousa",  "Full-Stack Developer",    "Candidatura Espontanea",
     [".NET","Angular","Azure"],     "Sugestao: Match com .NET Developer ou React Dev",  ACCENT_GREEN),
]):
    cx = Inches(0.5 + (i%2)*7.7); cy = Inches(2.2 + (i//2)*2.8)
    add_rect(slide, cx, cy, Inches(7.2), Inches(2.5), CARD_BG, radius=0.04)
    add_rect(slide, cx, cy, Inches(7.2), Inches(0.05), color)
    initials = "".join([w[0] for w in name.split()])
    add_icon_circle(slide, cx+Inches(0.3), cy+Inches(0.3), Inches(0.7), color, initials)
    add_text(slide, cx+Inches(1.2), cy+Inches(0.3),  Inches(4), Inches(0.3), name,       font_size=14, color=WHITE, bold=True)
    add_text(slide, cx+Inches(1.2), cy+Inches(0.6),  Inches(4), Inches(0.25), role,      font_size=11, color=MED_GRAY)
    add_text(slide, cx+Inches(1.2), cy+Inches(0.85), Inches(4), Inches(0.25), score_info, font_size=10, color=color)
    for j, skill in enumerate(skills):
        sx = cx+Inches(0.3+j*1.3)
        add_rect(slide, sx, cy+Inches(1.25), Inches(1.1), Inches(0.3), darken(color), radius=0.15)
        add_text(slide, sx, cy+Inches(1.27), Inches(1.1), Inches(0.28),
                 skill, font_size=9, color=color, alignment=PP_ALIGN.CENTER)
    add_rect(slide, cx+Inches(0.3), cy+Inches(1.75), Inches(6.5), Inches(0.5), DARK_BG2, radius=0.04)
    add_text(slide, cx+Inches(0.5), cy+Inches(1.8), Inches(6), Inches(0.4),
             suggestion, font_size=10, color=ACCENT_GREEN)

add_rect(slide, Inches(0.5), Inches(7.8), Inches(15), Inches(0.6), DARK_BG2, radius=0.03)
add_text(slide, Inches(1), Inches(7.85), Inches(14), Inches(0.5),
         "O Talent Pool garante que nenhum candidato qualificado se perde. "
         "Sugestoes automaticas de vagas futuras maximizam o retorno do pipeline.",
         font_size=12, color=LIGHT_GRAY)

add_bottom_bar(slide); add_page_number(slide, 11, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — DUPLICATE MERGE UI  (v2: UI interactiva real)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "QUALIDADE DE DADOS", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Merge Interactivo de Duplicados — Nova Feature v2",
         font_size=28, color=WHITE, bold=True, font_name="Segoe UI Light")

# ─── Panel esquerdo: Fila de duplicados (novo UI) ─────────────────
add_rect(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(6.7), CARD_BG, radius=0.04)
add_text(slide, Inches(0.8), Inches(1.65), Inches(5), Inches(0.4),
         "Fila de Revisao de Duplicados", font_size=15, color=WHITE, bold=True)
add_rect(slide, Inches(5.8), Inches(1.62), Inches(1.0), Inches(0.32),
         RGBColor(0x3D,0x2D,0x00), radius=0.12)
add_text(slide, Inches(5.8), Inches(1.64), Inches(1.0), Inches(0.3),
         "3 grupos", font_size=9, color=ACCENT_AMBER, bold=True, alignment=PP_ALIGN.CENTER)

# Grupo duplicado expandido
add_rect(slide, Inches(0.7), Inches(2.15), Inches(8.6), Inches(0.55),
         RGBColor(0x2A,0x1E,0x00), radius=0.04)
add_rect(slide, Inches(0.7), Inches(2.15), Inches(0.06), Inches(0.55), ACCENT_AMBER)
add_text(slide, Inches(0.9), Inches(2.2), Inches(5), Inches(0.3),
         "! 2 registos duplicados detectados  —  Telefone coincidente",
         font_size=11, color=ACCENT_AMBER, bold=True)
# Pill confiança
add_rect(slide, Inches(7.5), Inches(2.23), Inches(1.2), Inches(0.28),
         RGBColor(0x3D,0x1A,0x00), radius=0.12)
add_text(slide, Inches(7.5), Inches(2.25), Inches(1.2), Inches(0.26),
         "92% match", font_size=9, color=ACCENT_ROSE, bold=True, alignment=PP_ALIGN.CENTER)

# Corpo expandido: card primário + ícone merge + card duplicado
# Card primário
add_rect(slide, Inches(0.7), Inches(2.85), Inches(3.5), Inches(2.3),
         RGBColor(0x0F,0x2A,0x45), radius=0.05)
add_rect(slide, Inches(0.7), Inches(2.85), Inches(3.5), Inches(0.06),
         ACCENT_BLUE)
add_rect(slide, Inches(5.5), Inches(2.85), Inches(1.0), Inches(0.35),
         RGBColor(0x0D,0x27,0x1E), radius=0.12)
add_text(slide, Inches(5.5), Inches(2.87), Inches(1.0), Inches(0.33),
         "Primary", font_size=9, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.0), Inches(3.0), Inches(3), Inches(0.3),
         "Joao Santos", font_size=13, color=WHITE, bold=True)
add_text(slide, Inches(1.0), Inches(3.3), Inches(3), Inches(0.25),
         "joao.santos@email.com", font_size=10, color=LIGHT_GRAY)
add_text(slide, Inches(1.0), Inches(3.55), Inches(3), Inches(0.25),
         "+351 912 000 000", font_size=10, color=LIGHT_GRAY)
add_rect(slide, Inches(1.0), Inches(3.9), Inches(1.4), Inches(0.3),
         RGBColor(0x0D,0x27,0x1E), radius=0.08)
add_text(slide, Inches(1.0), Inches(3.92), Inches(1.4), Inches(0.28),
         "Score: 72", font_size=10, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(2.5), Inches(3.9), Inches(1.5), Inches(0.3),
         RGBColor(0x0B,0x19,0x35), radius=0.08)
add_text(slide, Inches(2.5), Inches(3.92), Inches(1.5), Inches(0.28),
         "SHORTLISTED", font_size=9, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Ícone merge central
add_icon_circle(slide, Inches(4.35), Inches(3.65), Inches(0.5), ACCENT_GREEN, "><")

# Card duplicado
add_rect(slide, Inches(4.95), Inches(2.85), Inches(3.5), Inches(2.3),
         CARD_BG, radius=0.05)
add_rect(slide, Inches(4.95), Inches(2.85), Inches(3.5), Inches(0.06), DARK_BG2)
add_text(slide, Inches(5.2), Inches(3.0), Inches(3), Inches(0.3),
         "J. Santos", font_size=13, color=WHITE, bold=True)
add_text(slide, Inches(5.2), Inches(3.3), Inches(3), Inches(0.25),
         "joao.santos@email.com", font_size=10, color=LIGHT_GRAY)
add_text(slide, Inches(5.2), Inches(3.55), Inches(3), Inches(0.25),
         "+351 912 000 000", font_size=10, color=LIGHT_GRAY)
add_rect(slide, Inches(5.2), Inches(3.9), Inches(1.4), Inches(0.3),
         RGBColor(0x2D,0x1C,0x00), radius=0.08)
add_text(slide, Inches(5.2), Inches(3.92), Inches(1.4), Inches(0.28),
         "Score: 65", font_size=10, color=ACCENT_AMBER, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(6.7), Inches(3.9), Inches(1.5), Inches(0.3),
         darken(MED_GRAY, 6), radius=0.08)
add_text(slide, Inches(6.7), Inches(3.92), Inches(1.5), Inches(0.28),
         "NEW", font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
# Botão "Set primary"
add_rect(slide, Inches(5.2), Inches(4.35), Inches(2.3), Inches(0.35), DARK_BG2, radius=0.08)
add_text(slide, Inches(5.2), Inches(4.37), Inches(2.3), Inches(0.33),
         "Definir como primario", font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Botão Merge
add_rect(slide, Inches(0.8), Inches(5.35), Inches(4.5), Inches(0.5), ACCENT_GREEN, radius=0.1)
add_text(slide, Inches(0.8), Inches(5.38), Inches(4.5), Inches(0.45),
         "Fundir em Joao Santos (primario)", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.8), Inches(5.92), Inches(7), Inches(0.3),
         "O registo duplicado e eliminado. Todos os dados movem para o primario.",
         font_size=9, color=MED_GRAY)

# Grupos colapsados
for i, (names, confidence) in enumerate([
    ("Ana R. / A. Rodrigues", "85% match"),
    ("Sofia M. / S. Martins",  "78% match"),
]):
    gy = Inches(6.45 + i*0.48)
    add_rect(slide, Inches(0.7), gy, Inches(8.6), Inches(0.4), RGBColor(0x1F,0x17,0x05), radius=0.04)
    add_rect(slide, Inches(0.7), gy, Inches(0.06), Inches(0.4), ACCENT_AMBER)
    add_text(slide, Inches(0.9), gy+Inches(0.06), Inches(5.5), Inches(0.28),
             f"! {names}", font_size=10, color=ACCENT_AMBER)
    add_text(slide, Inches(7.8), gy+Inches(0.06), Inches(1.3), Inches(0.28),
             confidence, font_size=9, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)

# ─── Panel direito: Exportacao & Auditoria ─────────────────────────
add_rect(slide, Inches(9.8), Inches(1.5), Inches(5.7), Inches(6.7), CARD_BG, radius=0.04)
add_text(slide, Inches(10.1), Inches(1.65), Inches(5), Inches(0.4),
         "Relatorios e Exportacoes", font_size=15, color=WHITE, bold=True)

for i, (name, desc, fmt, color) in enumerate([
    ("Exportacao Excel",   "Lista completa com scores,\nstatus e dados extraidos",  "XLSX", ACCENT_GREEN),
    ("Relatorio por Vaga", "Ranking de candidatos e\nmetricas de cada vaga",        "XLSX", ACCENT_BLUE),
    ("Audit Log",          "Registo de todas as accoes\nsensitivas no sistema",     "LOG",  ACCENT_AMBER),
    ("Merge History",      "Historico de duplicados\nfundidos com actor e data",    "LOG",  ACCENT_ROSE),
]):
    ey = Inches(2.3 + i*1.2)
    add_rect(slide, Inches(10.1), ey, Inches(5.1), Inches(1.0), DARK_BG2, radius=0.03)
    add_rect(slide, Inches(10.1), ey, Inches(0.08), Inches(1.0), color)
    add_text(slide, Inches(10.4), ey+Inches(0.1), Inches(3.5), Inches(0.3), name, font_size=13, color=WHITE, bold=True)
    add_text(slide, Inches(10.4), ey+Inches(0.4), Inches(3.5), Inches(0.5), desc, font_size=10, color=MED_GRAY)
    add_rect(slide, Inches(14.2), ey+Inches(0.3), Inches(0.8), Inches(0.35), color, radius=0.1)
    add_text(slide, Inches(14.2), ey+Inches(0.32), Inches(0.8), Inches(0.33),
             fmt, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide); add_page_number(slide, 12, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — INTEGRACOES
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "INTEGRACOES", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Ecossistema Microsoft 365 e Mais",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

for i, (title, desc, color, icon) in enumerate([
    ("Microsoft Graph\nMail",     "Sincronizacao bidirecional\nde emails com caixa de correio\ncareers@itsector.pt\n\nMail.Read | Mail.Send",                  ACCENT_BLUE,   "M"),
    ("Microsoft Graph\nCalendar", "Agendamento de entrevistas\ndirectamente no M365\n\nCalendars.ReadWrite\nConvites automaticos",                              ACCENT_PURPLE, "C"),
    ("Tesseract.js\nOCR",         "Reconhecimento de texto\nem CVs digitalizados\n\nMultiplos idiomas\nFallback automatico",                                     ACCENT_CYAN,   "O"),
    ("ExcelJS\nExportacao",       "Geracao de relatorios\nExcel com dados completos\nde candidatos e metricas\n\nExportacao local",                              ACCENT_GREEN,  "X"),
]):
    cx = Inches(0.5 + i*3.85); cy = Inches(1.8)
    add_rect(slide, cx, cy, Inches(3.6), Inches(4.5), CARD_BG, radius=0.04)
    add_rect(slide, cx, cy, Inches(3.6), Inches(0.06), color)
    add_icon_circle(slide, cx+Inches(1.35), cy+Inches(0.4), Inches(0.85), color, icon)
    add_text(slide, cx+Inches(0.3), cy+Inches(1.5), Inches(3), Inches(0.7),
             title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, cx+Inches(0.3), cy+Inches(2.3), Inches(3), Inches(2.0),
             desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.5), Inches(6.7), Inches(15), Inches(1.5), CARD_BG, radius=0.03)
add_text(slide, Inches(1), Inches(6.85), Inches(5), Inches(0.4),
         "Inteligencia Artificial (Opcional)", font_size=16, color=ACCENT_PURPLE, bold=True)
add_text(slide, Inches(1), Inches(7.3), Inches(13), Inches(0.7),
         "Suporte para OpenAI, Azure OpenAI e Anthropic — para resumos de CV, sugestoes de email e classificacao.\n"
         "Totalmente opcional: o sistema funciona com heuristicas locais sem API externa.",
         font_size=12, color=LIGHT_GRAY)
for i, (name, color) in enumerate([("OpenAI",ACCENT_GREEN),("Azure OpenAI",ACCENT_BLUE),("Anthropic",ACCENT_PURPLE)]):
    add_rect(slide, Inches(11+i*1.5), Inches(6.95), Inches(1.3), Inches(0.35), color, radius=0.1)
    add_text(slide, Inches(11+i*1.5), Inches(6.97), Inches(1.3), Inches(0.33),
             name, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide); add_page_number(slide, 13, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — TECNOLOGIA  (v2.3 — novo slide)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "TECNOLOGIA", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(14), Inches(0.6),
         "O Que Usamos e Como Funciona — Stack Sem Vendor Lock-in",
         font_size=28, color=WHITE, bold=True, font_name="Segoe UI Light")

tech_cards = [
    ("Frontend", [
        "Next.js 15 App Router",
        "React 19 (Server Components)",
        "TypeScript  |  Tailwind CSS",
        "TanStack Table  |  Recharts",
    ], ACCENT_BLUE, "F"),
    ("Backend & API", [
        "Next.js API Routes",
        "Services + Repository Layer",
        "Zod Validation em todas as APIs",
        "Rate Limiting 5 req / 15 min",
    ], ACCENT_CYAN, "B"),
    ("Base de Dados", [
        "Prisma 6 ORM",
        "SQLite (dev)  ->  PostgreSQL (prod)",
        "14+ modelos, migrations automaticas",
        "Indices optimizados para pesquisa",
    ], ACCENT_GREEN, "D"),
    ("Microsoft 365", [
        "Microsoft Graph API",
        "Mail.Read  |  Mail.Send  |  Mail.ReadWrite",
        "Calendars.ReadWrite",
        "OAuth 2.0 Client Credentials",
    ], ACCENT_PURPLE, "M"),
    ("IA — Opcional", [
        "Azure OpenAI  |  OpenAI GPT-4",
        "Anthropic Claude",
        "Sem IA: heuristicas locais activas",
        "Zero dependencia de API externa",
    ], ACCENT_AMBER, "AI"),
    ("Seguranca", [
        "HMAC-signed session cookies",
        "requireAdmin  |  requireUser guards",
        "Audit log com actor e timestamp",
        "isActive guard + Retry-After 429",
    ], ACCENT_ROSE, "S"),
]

for i, (title, items, color, icon) in enumerate(tech_cards):
    row = i // 3; col = i % 3
    cx = Inches(0.5 + col * 5.1); cy = Inches(1.7 + row * 3.1)
    add_rect(slide, cx, cy, Inches(4.8), Inches(2.8), CARD_BG, radius=0.04)
    add_rect(slide, cx, cy, Inches(4.8), Inches(0.06), color)
    add_icon_circle(slide, cx + Inches(0.3), cy + Inches(0.25), Inches(0.6), color, icon)
    add_text(slide, cx + Inches(1.1), cy + Inches(0.28), Inches(3.5), Inches(0.4),
             title, font_size=15, color=WHITE, bold=True)
    for j, item in enumerate(items):
        add_text(slide, cx + Inches(0.5), cy + Inches(0.95 + j * 0.44), Inches(4.1), Inches(0.38),
                 item, font_size=11, color=LIGHT_GRAY)

add_rect(slide, Inches(0.5), Inches(8.0), Inches(15), Inches(0.65), DARK_BG2, radius=0.03)
add_text(slide, Inches(0.8), Inches(8.08), Inches(14.2), Inches(0.5),
         "Corre em qualquer servidor Windows / Linux com Node.js 20+  —  Sem dependencia de cloud obrigatoria.  "
         "Azure App Service / Docker / On-premises: pronto para qualquer ambiente.",
         font_size=11, color=LIGHT_GRAY)

add_bottom_bar(slide); add_page_number(slide, 14, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 15 — ARQUITECTURA  (v2: camada de Seguranca destacada)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "ARQUITECTURA", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "Stack Moderna, Segura e Escalavel",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

layers = [
    ("FRONTEND",       Inches(1.55), [
        ("Next.js 15",        ACCENT_BLUE),
        ("React 19",          ACCENT_CYAN),
        ("Tailwind CSS",      ACCENT_GREEN),
        ("TanStack Table",    ACCENT_PURPLE),
        ("Pagination UI",     ACCENT_AMBER),
        ("Recharts",          ACCENT_ROSE),
    ]),
    ("BACKEND / API",  Inches(3.05), [
        ("API Routes",        ACCENT_BLUE),
        ("Services Layer",    ACCENT_CYAN),
        ("Repository Pattern",ACCENT_GREEN),
        ("TypeScript",        ACCENT_PURPLE),
        ("Zod Validation",    ACCENT_AMBER),
        ("Rate Limiting",     ACCENT_ROSE),
    ]),
    ("SEGURANCA",      Inches(4.55), [
        ("HMAC Sessions",     ACCENT_BLUE),
        ("requireAdmin",      ACCENT_CYAN),
        ("requireUser",       ACCENT_GREEN),
        ("5 req/15min Limit", ACCENT_PURPLE),
        ("Audit Logging",     ACCENT_AMBER),
        ("isActive Guard",    ACCENT_ROSE),
    ]),
    ("SERVICOS",       Inches(6.05), [
        ("Inbox Sync",        ACCENT_BLUE),
        ("CV Parsing",        ACCENT_CYAN),
        ("Scoring Engine",    ACCENT_GREEN),
        ("Duplicate Merge",   ACCENT_PURPLE),
        ("Email Service",     ACCENT_AMBER),
        ("Calendar Service",  ACCENT_ROSE),
    ]),
    ("DADOS & STORAGE",Inches(7.55), [
        ("SQLite + Prisma",   ACCENT_BLUE),
        ("Filesystem Local",  ACCENT_CYAN),
        ("14+ Modelos DB",    ACCENT_GREEN),
        ("DB Indexes",        ACCENT_PURPLE),
        ("Migrations",        ACCENT_AMBER),
        ("Seed Data",         ACCENT_ROSE),
    ]),
]

for layer_name, y, items in layers:
    is_security = layer_name == "SEGURANCA"
    bg_color = RGBColor(0x1A, 0x1F, 0x2E) if not is_security else RGBColor(0x14, 0x22, 0x14)
    add_rect(slide, Inches(0.5), y, Inches(15), Inches(1.25), bg_color, radius=0.03)
    if is_security:
        add_rect(slide, Inches(0.5), y, Inches(0.06), Inches(1.25), ACCENT_GREEN)
    label_color = ACCENT_GREEN if is_security else ACCENT_BLUE
    add_text(slide, Inches(0.8), y+Inches(0.1), Inches(2.3), Inches(0.35),
             layer_name, font_size=12, color=label_color, bold=True)
    for i, (tech, color) in enumerate(items):
        tx = Inches(0.8 + i*2.45)
        add_rect(slide, tx, y+Inches(0.55), Inches(2.2), Inches(0.5), darken(color,5), radius=0.08)
        add_rect(slide, tx, y+Inches(0.55), Inches(0.06), Inches(0.5), color)
        add_text(slide, tx+Inches(0.2), y+Inches(0.62), Inches(1.9), Inches(0.35),
                 tech, font_size=11, color=color, bold=True)
    if y < Inches(7.5):
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(7.7), y+Inches(1.25), Inches(0.3), Inches(0.2))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT_CYAN; arrow.line.fill.background()

add_bottom_bar(slide); add_page_number(slide, 15, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 16 — ROADMAP  (v2: Fases 1-4 CONCLUIDO, Fase 5 PLANEADO)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
         "ROADMAP", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1), Inches(0.7), Inches(12), Inches(0.6),
         "4 Fases Concluidas — Fase 5 em Planeamento",
         font_size=30, color=WHITE, bold=True, font_name="Segoe UI Light")

add_rect(slide, Inches(1), Inches(2.7), Inches(14), Inches(0.06), ACCENT_BLUE)

phases = [
    ("FASE 1", "Fundacao ATS",              "CONCLUIDO",  ACCENT_GREEN,
     ["Auth + sessoes HMAC", "Dashboard KPIs reais", "Sync M365 Mail",
      "Parsing PDF/DOCX", "Motor de scoring", "Templates + envio Graph"]),
    ("FASE 2", "Produtividade Operacional", "CONCLUIDO",  ACCENT_GREEN,
     ["Importacao JSON vagas", "Bulk actions auditados", "Exportacao Excel",
      "Indices DB (14+)", "Pipeline view", "Talent Pool matching"]),
    ("FASE 3", "UX & Admin",               "CONCLUIDO",  ACCENT_GREEN,
     ["Timeline unificada", "Notas estruturadas", "Admin users toggle",
      "Role guards (req/admin)", "Candidate detail completo", "Audit log actor"]),
    ("FASE 4", "Seguranca & Qualidade",    "CONCLUIDO",  ACCENT_GREEN,
     ["Rate limiting login", "Pagination controls URL", "Duplicate merge UI",
      "Merge API hardened", "Retry-After headers", "isActive guard"]),
    ("FASE 5", "Cobertura Completa",       "PLANEADO",   ACCENT_AMBER,
     ["CV upload manual", "Template CRUD UI", "Relatorios charts",
      "Auth em todas as APIs", "Calendar interview UI", "Settings editor"]),
]

# Linha de tempo com 5 fases — mais compacta
for i, (phase, title, status, color, items) in enumerate(phases):
    cx = Inches(0.3 + i*3.1)
    add_icon_circle(slide, cx+Inches(1.1), Inches(2.45), Inches(0.5), color, str(i+1))
    card = add_rect(slide, cx, Inches(3.1), Inches(2.9), Inches(5.1), CARD_BG, radius=0.04)
    add_rect(slide, cx, Inches(3.1), Inches(2.9), Inches(0.05), color)
    add_text(slide, cx+Inches(0.2), Inches(3.3),  Inches(1.1), Inches(0.28), phase,  font_size=11, color=color, bold=True)
    add_rect(slide, cx+Inches(1.4), Inches(3.28), Inches(1.3), Inches(0.32), darken(color), radius=0.12)
    add_text(slide, cx+Inches(1.4), Inches(3.30), Inches(1.3), Inches(0.30),
             status, font_size=8, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, cx+Inches(0.2), Inches(3.72), Inches(2.5), Inches(0.38),
             title, font_size=13, color=WHITE, bold=True)
    for j, item in enumerate(items):
        iy = Inches(4.2 + j*0.46)
        add_rect(slide, cx+Inches(0.3), iy+Inches(0.05), Inches(0.12), Inches(0.12), color, radius=0.5)
        add_text(slide, cx+Inches(0.55), iy, Inches(2.2), Inches(0.3),
                 item, font_size=10, color=LIGHT_GRAY)

add_bottom_bar(slide); add_page_number(slide, 16, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 17 — CTA  (v2: stats actualizadas)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-3), Inches(12), Inches(12))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x15,0x20,0x3C); c1.line.fill.background()
c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(3), Inches(10), Inches(10))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x1E,0x29,0x3B); c2.line.fill.background()

add_gradient_bar(slide, Inches(0), Inches(0), Inches(16), Inches(0.08), ACCENT_BLUE, ACCENT_CYAN)

add_text(slide, Inches(1), Inches(2.0), Inches(14), Inches(1.0),
         "Pronto para Transformar", font_size=48, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")
add_text(slide, Inches(1), Inches(3.0), Inches(14), Inches(1.0),
         "o Recrutamento?", font_size=48, color=ACCENT_CYAN, bold=True,
         alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")

add_rect(slide, Inches(6.5), Inches(4.2), Inches(3), Inches(0.04), ACCENT_BLUE)
add_text(slide, Inches(2), Inches(4.6), Inches(12), Inches(0.8),
         "O ITSector Talent Inbox ATS v2 esta pronto para implementacao imediata.\n"
         "4 Fases completas. 100% local. Seguro. Inteligente.",
         font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

for i, (val, label) in enumerate([
    ("19+",  "Funcionalidades\nCompletas"),
    ("14+",  "Modelos de\nDados"),
    ("4",    "Fases\nConcluidas"),
    ("100%", "Local\n& Seguro"),
]):
    sx = Inches(1.5 + i*3.5)
    add_text(slide, sx, Inches(5.8), Inches(2.5), Inches(0.6),
             val, font_size=42, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, sx, Inches(6.4), Inches(2.5), Inches(0.6),
             label, font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(7.3), Inches(5), Inches(0.7), ACCENT_BLUE, radius=0.1)
add_text(slide, Inches(5.5), Inches(7.35), Inches(5), Inches(0.6),
         "Comecar Agora  >>", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(3), Inches(8.2), Inches(10), Inches(0.4),
         "ITSector  |  careers@itsector.pt  |  Talent Inbox ATS  v2.0",
         font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide); add_page_number(slide, 17, TOTAL_SLIDES)


# ─── GUARDAR ─────────────────────────────────────────────────────
ROOT = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(ROOT, "ITSector_Talent_ATS_Comercial_v2.3_PT.pptx")
prs.save(output_path)
print(f"Apresentacao guardada em: {output_path}")
print(f"Total de slides: {TOTAL_SLIDES}")
